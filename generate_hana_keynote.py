#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
하나은행 관통 장표 2장 — "AI 시대 금융의 통제권 전쟁"
slide 1: Why Now × Why KT — 통제권 전쟁의 답
slide 2: 통제 × 운영 — AI는 구축이 아니라 운영에서 무너진다
"""

from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─── 팔레트 ──────────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x07, 0x10, 0x1E)   # 더 어두운 네이비 (긴장감)
NAVY2       = RGBColor(0x0D, 0x1F, 0x3C)
NAVY3       = RGBColor(0x13, 0x2A, 0x4D)
NAVY4       = RGBColor(0x18, 0x33, 0x5A)
CYAN        = RGBColor(0x00, 0xD4, 0xFF)
CYAN_DIM    = RGBColor(0x00, 0x88, 0xA8)
ORANGE      = RGBColor(0xFF, 0x6B, 0x35)
YELLOW      = RGBColor(0xFF, 0xD7, 0x00)
GOLD        = RGBColor(0xE5, 0xB2, 0x2D)
GREEN       = RGBColor(0x00, 0xE0, 0x96)
PURPLE      = RGBColor(0xB4, 0x6B, 0xFF)
RED         = RGBColor(0xFF, 0x4D, 0x4D)
RED_DEEP    = RGBColor(0xB8, 0x2A, 0x2A)
KT_RED      = RGBColor(0xE3, 0x14, 0x35)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xD6, 0xE0)
MID_BLUE    = RGBColor(0x1A, 0x4A, 0x7A)
DARK_TEXT   = RGBColor(0xA0, 0xB8, 0xCC)

FONT_KO = "맑은 고딕"
SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)
FOOTER_TEXT = "하나은행 비정형 데이터 자산화 플랫폼 제안  |  2026.05  |  KT B2B 수주전략팀"


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=None, shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    if line_color:
        sh.line.color.rgb = line_color
        if line_width:
            sh.line.width = line_width
    else:
        sh.line.fill.background()
    return sh


def add_textbox(slide, text, x, y, w, h,
                font_size=12, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.08)
    tf.margin_right = Cm(0.08)
    tf.margin_top = Cm(0.04)
    tf.margin_bottom = Cm(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_KO
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_rich_textbox(slide, runs, x, y, w, h, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.05)
    tf.margin_bottom = Cm(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    for (text, fs, bold, color) in runs:
        r = p.add_run()
        r.text = text
        r.font.name = FONT_KO
        r.font.size = Pt(fs)
        r.font.bold = bold
        r.font.color.rgb = color
    return txb


def add_chevron(slide, x, y, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def add_footer(slide, page_num, total=3):
    add_rect(slide, Cm(1.5), SLIDE_H - Cm(1.0),
             SLIDE_W - Cm(3), Cm(0.04), MID_BLUE)
    add_textbox(slide, FOOTER_TEXT,
                Cm(1.5), SLIDE_H - Cm(0.9),
                SLIDE_W - Cm(5), Cm(0.7),
                font_size=8, color=DARK_TEXT, align=PP_ALIGN.LEFT)
    add_textbox(slide, f"{page_num} / {total}",
                SLIDE_W - Cm(3.5), SLIDE_H - Cm(0.9),
                Cm(2), Cm(0.7),
                font_size=8, color=DARK_TEXT, align=PP_ALIGN.RIGHT)


# ─── 슬라이드 1: 4계층 토대 구조 — KT 인프라 위에 쌓인 안정감 ───────────────────

def slide_01_war(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)

    # ── 상단 헤더 띠 ──────────────────────────────────────────────────────
    add_rect(sl, 0, 0, SLIDE_W, Cm(2.9), NAVY2)
    add_rect(sl, 0, 0, Cm(0.25), Cm(2.9), GOLD)

    # 태그
    add_rect(sl, Cm(1.2), Cm(0.35), Cm(9.5), Cm(0.5), GOLD)
    add_textbox(sl, "관통 메시지  ·  THE FOUNDATION OF TRUST",
                Cm(1.2), Cm(0.35), Cm(9.5), Cm(0.5),
                font_size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # 메인 헤드라인
    add_textbox(sl,
                "리스크를 통제할 때, 고객과 직원의 안정감이 지속됩니다.",
                Cm(1.2), Cm(0.95), SLIDE_W - Cm(2.4), Cm(1.05),
                font_size=23, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # 서브
    add_rich_textbox(sl, [
        ("수집부터 가공까지  ─  ", 12, False, LIGHT_GRAY),
        ("AI 통제 가능한 비정형 데이터 플랫폼", 12, True, CYAN),
        ("이 그 토대입니다.", 12, False, LIGHT_GRAY),
    ], Cm(1.2), Cm(2.1), SLIDE_W - Cm(2.4), Cm(0.6))

    # 우측 상단 안내선
    add_textbox(sl, "▼  표면 (사람이 느끼는 것)",
                SLIDE_W - Cm(8), Cm(0.95), Cm(7), Cm(0.5),
                font_size=9, italic=True, color=DARK_TEXT, align=PP_ALIGN.RIGHT)
    add_textbox(sl, "▲  토대 (받쳐주는 것)",
                SLIDE_W - Cm(8), Cm(1.5), Cm(7), Cm(0.5),
                font_size=9, italic=True, color=GOLD, align=PP_ALIGN.RIGHT)

    # ── 본문 4계층 영역 ───────────────────────────────────────────────────
    body_x = Cm(1.2)
    body_w = SLIDE_W - Cm(2.4)
    label_w = Cm(2.8)
    layer_x = body_x + label_w + Cm(0.25)
    layer_w = body_w - label_w - Cm(0.25)

    # ─ LAYER 04 (TOP) ─ 안정감의 지속 ────────────────────────────────
    L1_y = Cm(3.1)
    L1_h = Cm(2.7)
    # 따스한 글로우 배경
    add_rect(sl, layer_x, L1_y, layer_w, L1_h, RGBColor(0x1A, 0x17, 0x0C),
             GOLD, Pt(0.6))
    add_rect(sl, layer_x, L1_y, Cm(0.18), L1_h, GOLD)
    # 좌측 라벨
    add_rect(sl, body_x, L1_y, label_w, L1_h, NAVY2, GOLD, Pt(0.4))
    add_textbox(sl, "LAYER 04",
                body_x, L1_y + Cm(0.4), label_w, Cm(0.45),
                font_size=10, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(sl, "표면",
                body_x, L1_y + Cm(0.85), label_w, Cm(0.7),
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, "Surface",
                body_x, L1_y + Cm(1.55), label_w, Cm(0.4),
                font_size=8, italic=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    # 내부 헤더
    add_textbox(sl, "안정감의 지속  ─  사람이 흔들리지 않는다",
                layer_x + Cm(0.5), L1_y + Cm(0.2), layer_w - Cm(1), Cm(0.55),
                font_size=12, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    # 2 카드: 대고객 / 대직원
    card_y = L1_y + Cm(0.85)
    card_h = L1_h - Cm(1.05)
    card_gap = Cm(0.4)
    card_w = (layer_w - Cm(1.0) - card_gap) / 2

    # 대고객
    cx1 = layer_x + Cm(0.5)
    add_rect(sl, cx1, card_y, card_w, card_h, NAVY3, CYAN, Pt(0.6))
    # 아이콘 자리 (오벌)
    add_rect(sl, cx1 + Cm(0.4), card_y + Cm(0.35), Cm(0.85), Cm(0.85),
             CYAN, shape=MSO_SHAPE.OVAL)
    add_textbox(sl, "고",
                cx1 + Cm(0.4), card_y + Cm(0.35), Cm(0.85), Cm(0.85),
                font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(sl, "대고객",
                cx1 + Cm(1.45), card_y + Cm(0.3), card_w - Cm(1.7), Cm(0.5),
                font_size=11, bold=True, color=CYAN, align=PP_ALIGN.LEFT)
    add_textbox(sl, "흔들리지 않는 응대",
                cx1 + Cm(1.45), card_y + Cm(0.78), card_w - Cm(1.7), Cm(0.65),
                font_size=15, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # 대직원
    cx2 = cx1 + card_w + card_gap
    add_rect(sl, cx2, card_y, card_w, card_h, NAVY3, YELLOW, Pt(0.6))
    add_rect(sl, cx2 + Cm(0.4), card_y + Cm(0.35), Cm(0.85), Cm(0.85),
             YELLOW, shape=MSO_SHAPE.OVAL)
    add_textbox(sl, "직",
                cx2 + Cm(0.4), card_y + Cm(0.35), Cm(0.85), Cm(0.85),
                font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(sl, "대직원",
                cx2 + Cm(1.45), card_y + Cm(0.3), card_w - Cm(1.7), Cm(0.5),
                font_size=11, bold=True, color=YELLOW, align=PP_ALIGN.LEFT)
    add_textbox(sl, "흔들리지 않는 판단",
                cx2 + Cm(1.45), card_y + Cm(0.78), card_w - Cm(1.7), Cm(0.65),
                font_size=15, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # ↑ 연결 (보호받는다)
    a1_y = L1_y + L1_h
    add_textbox(sl, "▲   보호받는다",
                layer_x, a1_y, layer_w, Cm(0.4),
                font_size=9, italic=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    # ─ LAYER 03 ─ 리스크 통제 (방패) ─────────────────────────────────
    L2_y = a1_y + Cm(0.45)
    L2_h = Cm(2.1)
    add_rect(sl, layer_x, L2_y, layer_w, L2_h, NAVY2, RED, Pt(0.7))
    add_rect(sl, layer_x, L2_y, Cm(0.18), L2_h, RED)
    add_rect(sl, body_x, L2_y, label_w, L2_h, NAVY2, RED, Pt(0.4))
    add_textbox(sl, "LAYER 03",
                body_x, L2_y + Cm(0.25), label_w, Cm(0.45),
                font_size=10, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_textbox(sl, "통제",
                body_x, L2_y + Cm(0.7), label_w, Cm(0.7),
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, "Control",
                body_x, L2_y + Cm(1.4), label_w, Cm(0.4),
                font_size=8, italic=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    add_textbox(sl, "리스크 통제  ─  4축 가드레일",
                layer_x + Cm(0.5), L2_y + Cm(0.15), layer_w - Cm(1), Cm(0.5),
                font_size=12, bold=True, color=RED, align=PP_ALIGN.LEFT)

    # 4축 가드레일 (방패 모양 이미지)
    axes = [
        ("출처", "Provenance", YELLOW),
        ("감사", "Audit",      PURPLE),
        ("접근", "Access",     CYAN),
        ("품질", "Quality",    GREEN),
    ]
    ax_y = L2_y + Cm(0.75)
    ax_h = L2_h - Cm(0.95)
    ax_gap = Cm(0.25)
    ax_w = (layer_w - Cm(1.0) - ax_gap * 3) / 4
    for i, (name, en, color) in enumerate(axes):
        ax = layer_x + Cm(0.5) + i * (ax_w + ax_gap)
        add_rect(sl, ax, ax_y, ax_w, ax_h, NAVY3, color, Pt(0.5),
                 shape=MSO_SHAPE.PENTAGON)
        d = Cm(0.32)
        add_rect(sl, ax + ax_w/2 - d/2, ax_y + Cm(0.15), d, d,
                 color, shape=MSO_SHAPE.OVAL)
        add_textbox(sl, name,
                    ax, ax_y + Cm(0.5), ax_w, Cm(0.55),
                    font_size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_textbox(sl, en,
                    ax, ax_y + Cm(1.0), ax_w, Cm(0.4),
                    font_size=8, italic=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # ↑ 연결
    a2_y = L2_y + L2_h
    add_textbox(sl, "▲   통제 가능한 형태로 가공되어 올라온다",
                layer_x, a2_y, layer_w, Cm(0.4),
                font_size=9, italic=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    # ─ LAYER 02 ─ 비정형 데이터 플랫폼 (수집→저장→가공) ──────────────
    L3_y = a2_y + Cm(0.45)
    L3_h = Cm(2.3)
    add_rect(sl, layer_x, L3_y, layer_w, L3_h, NAVY2, CYAN, Pt(0.7))
    add_rect(sl, layer_x, L3_y, Cm(0.18), L3_h, CYAN)
    add_rect(sl, body_x, L3_y, label_w, L3_h, NAVY2, CYAN, Pt(0.4))
    add_textbox(sl, "LAYER 02",
                body_x, L3_y + Cm(0.3), label_w, Cm(0.45),
                font_size=10, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    add_textbox(sl, "플랫폼",
                body_x, L3_y + Cm(0.78), label_w, Cm(0.7),
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, "Platform",
                body_x, L3_y + Cm(1.5), label_w, Cm(0.4),
                font_size=8, italic=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    add_textbox(sl, "비정형 데이터 플랫폼  ─  수집부터 가공까지 통제 가능한 흐름",
                layer_x + Cm(0.5), L3_y + Cm(0.15), layer_w - Cm(1), Cm(0.5),
                font_size=12, bold=True, color=CYAN, align=PP_ALIGN.LEFT)

    # 3단 chevron 흐름 (수집→저장→가공)
    flow_y = L3_y + Cm(0.75)
    flow_h = L3_h - Cm(0.95)
    stages = [
        ("01", "수집", "Ingest",  "원천 + 출처 태그"),
        ("02", "저장", "Govern",  "보존 + 접근 통제"),
        ("03", "가공", "Process", "AI Ready化 + 정합성"),
    ]
    fl_total = layer_w - Cm(1.0)
    fl_gap = Cm(0.2)
    fl_w = (fl_total - fl_gap * 2 - Cm(0.6)) / 3
    fl_x = layer_x + Cm(0.5)
    for i, (idx, ko, en, sub) in enumerate(stages):
        fx = fl_x + i * (fl_w + fl_gap)
        add_rect(sl, fx, flow_y, fl_w, flow_h, NAVY3, CYAN, Pt(0.5))
        # 헤더 바
        add_rect(sl, fx, flow_y, fl_w, Cm(0.55), CYAN_DIM)
        add_textbox(sl, idx,
                    fx + Cm(0.25), flow_y + Cm(0.08), Cm(1), Cm(0.4),
                    font_size=10, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
        add_textbox(sl, ko,
                    fx, flow_y + Cm(0.08), fl_w, Cm(0.4),
                    font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_textbox(sl, en,
                    fx, flow_y + Cm(0.65), fl_w, Cm(0.4),
                    font_size=9, italic=True, color=CYAN, align=PP_ALIGN.CENTER)
        add_textbox(sl, sub,
                    fx + Cm(0.25), flow_y + Cm(1.05), fl_w - Cm(0.5), Cm(0.5),
                    font_size=10, color=WHITE, align=PP_ALIGN.CENTER)
        if i < 2:
            add_chevron(sl, fx + fl_w + Cm(0.0), flow_y + flow_h/2 - Cm(0.22),
                        Cm(0.3), Cm(0.45), CYAN)

    # ↑ 연결
    a3_y = L3_y + L3_h
    add_textbox(sl, "▲   받쳐주고  ·  멈추지 않게 한다",
                layer_x, a3_y, layer_w, Cm(0.4),
                font_size=9, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

    # ─ LAYER 01 (BASE) ─ KT INFRASTRUCTURE ──────────────────────────
    L4_y = a3_y + Cm(0.45)
    L4_h = Cm(2.9)
    # 두꺼운 토대 BAR (KT_RED 강조)
    add_rect(sl, layer_x, L4_y, layer_w, L4_h, KT_RED)
    # 받침 그림자
    add_rect(sl, layer_x + Cm(0.3), L4_y + L4_h,
             layer_w - Cm(0.6), Cm(0.18), RGBColor(0x60, 0x08, 0x18))
    # 좌측 라벨
    add_rect(sl, body_x, L4_y, label_w, L4_h, KT_RED, WHITE, Pt(0.5))
    add_textbox(sl, "LAYER 01",
                body_x, L4_y + Cm(0.35), label_w, Cm(0.45),
                font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, "토대",
                body_x, L4_y + Cm(0.85), label_w, Cm(0.85),
                font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, "Foundation",
                body_x, L4_y + Cm(1.7), label_w, Cm(0.4),
                font_size=8, italic=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    # WHY KT 마크
    add_rect(sl, body_x + Cm(0.4), L4_y + L4_h - Cm(0.75), label_w - Cm(0.8), Cm(0.5),
             WHITE)
    add_textbox(sl, "WHY KT",
                body_x + Cm(0.4), L4_y + L4_h - Cm(0.75), label_w - Cm(0.8), Cm(0.5),
                font_size=10, bold=True, color=KT_RED, align=PP_ALIGN.CENTER)

    # 내부 핵심 카피
    add_textbox(sl, "INFRASTRUCTURE  ·  멈추지 않는 운영 DNA",
                layer_x + Cm(0.6), L4_y + Cm(0.25), layer_w - Cm(1.2), Cm(0.5),
                font_size=10, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # 가장 굵직한 한 줄
    add_textbox(sl,
                "통신을 멈추지 않게 운영해온 단 하나의 회사",
                layer_x + Cm(0.6), L4_y + Cm(0.8), layer_w - Cm(1.2), Cm(0.95),
                font_size=22, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_textbox(sl,
                "이 모든 레이어를 받칠 수 있는 토대는, KT의 인프라 운영 DNA 하나입니다.",
                layer_x + Cm(0.6), L4_y + Cm(1.85), layer_w - Cm(1.2), Cm(0.55),
                font_size=12, italic=True, color=WHITE, align=PP_ALIGN.LEFT)

    # 기둥 표현 (인프라가 받치는 시각적 단서)
    pillar_y = L4_y + L4_h - Cm(0.95)
    pillar_w = Cm(0.5)
    pillar_h = Cm(0.7)
    pillar_count = 8
    pillar_total_w = layer_w - Cm(1.2)
    pillar_gap = (pillar_total_w - pillar_w * pillar_count) / (pillar_count - 1)
    for i in range(pillar_count):
        px = layer_x + Cm(0.6) + i * (pillar_w + pillar_gap)
        add_rect(sl, px, pillar_y, pillar_w, pillar_h,
                 RGBColor(0xFF, 0xB0, 0xB8), shape=MSO_SHAPE.TRAPEZOID)

    add_footer(sl, 1)


# ─── 슬라이드 2: 통제 × 운영 (구축이 아니라 운영) ──────────────────────────────

def slide_02_operation(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)

    # ── 상단 헤더 띠 ──────────────────────────────────────────────────────
    add_rect(sl, 0, 0, SLIDE_W, Cm(3.4), NAVY2)
    add_rect(sl, 0, 0, Cm(0.25), Cm(3.4), GOLD)

    add_rect(sl, Cm(1.2), Cm(0.45), Cm(9.0), Cm(0.55), GOLD)
    add_textbox(sl, "통제 × 운영  ·  HOW WE WIN THE CONTROL WAR",
                Cm(1.2), Cm(0.45), Cm(9.0), Cm(0.55),
                font_size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # 메인 헤드라인
    add_textbox(sl,
                "AI는 구축이 아니라 운영에서 무너집니다.",
                Cm(1.2), Cm(1.1), SLIDE_W - Cm(2.4), Cm(1.2),
                font_size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # 서브
    add_rich_textbox(sl, [
        ("초기 구축보다 ", 14, False, LIGHT_GRAY),
        ("장기 운영비", 14, True, ORANGE),
        ("가 무섭다  ─  지속 가능한 ", 14, False, LIGHT_GRAY),
        ("통제 체계", 14, True, CYAN),
        ("만이 답이다.", 14, False, LIGHT_GRAY),
    ], Cm(1.2), Cm(2.5), SLIDE_W - Cm(2.4), Cm(0.7))

    # ── 상단: 통제 가능한 AI 파이프라인 ───────────────────────────────────
    section_label_y = Cm(3.85)
    add_rect(sl, Cm(1.2), section_label_y, Cm(0.4), Cm(0.55), CYAN)
    add_textbox(sl, "PART 1.  통제 — 출처가 박힌 데이터만 학습하고, 인용 가능한 응답만 전달한다",
                Cm(1.8), section_label_y, SLIDE_W - Cm(3), Cm(0.55),
                font_size=12, bold=True, color=CYAN, align=PP_ALIGN.LEFT)

    # 통제축 띠 (얇게)
    ctrl_y = section_label_y + Cm(0.7)
    ctrl_h = Cm(0.9)
    add_rect(sl, Cm(1.2), ctrl_y, SLIDE_W - Cm(2.4), ctrl_h, NAVY3, MID_BLUE, Pt(0.5))
    axes = [("출처", YELLOW), ("감사", PURPLE), ("접근", CYAN), ("품질", GREEN)]
    add_textbox(sl, "▼  Cross-cutting Controls",
                Cm(1.5), ctrl_y + Cm(0.18), Cm(5), Cm(0.55),
                font_size=9, italic=True, color=DARK_TEXT, align=PP_ALIGN.LEFT)
    ax_x = Cm(8.0)
    for name, color in axes:
        add_rect(sl, ax_x, ctrl_y + Cm(0.28), Cm(0.3), Cm(0.3), color, shape=MSO_SHAPE.OVAL)
        add_textbox(sl, name,
                    ax_x + Cm(0.4), ctrl_y + Cm(0.2), Cm(2.5), Cm(0.5),
                    font_size=11, bold=True, color=color, align=PP_ALIGN.LEFT)
        ax_x += Cm(3.2)

    # 4-stage 파이프라인
    pipe_y = ctrl_y + ctrl_h + Cm(0.25)
    pipe_h = Cm(3.4)
    diag_w = SLIDE_W - Cm(2.4)
    stages = [
        ("수집",  "Ingest",  "원천 → 출처 태그 부착",  "raw + meta",      CYAN),
        ("저장",  "Govern",  "출처 보존 + 접근 통제",  "stored + audit",  PURPLE),
        ("가공",  "Process", "AI Ready化 + 정합성",    "ai-ready",        GREEN),
        ("활용",  "Serve",   "출처 인용 + 응답 검증",   "served + cited",  ORANGE),
    ]
    gap = Cm(0.2)
    stage_w = (diag_w - 3 * gap) / 4
    for i, (ko, en, sub, state, color) in enumerate(stages):
        sx = Cm(1.2) + i * (stage_w + gap)
        add_rect(sl, sx, pipe_y, stage_w, pipe_h, NAVY2, color, Pt(1.0))
        # 상단 헤더 바
        add_rect(sl, sx, pipe_y, stage_w, Cm(1.1), color)
        add_textbox(sl, f"0{i+1}",
                    sx + Cm(0.25), pipe_y + Cm(0.1), Cm(0.7), Cm(0.4),
                    font_size=9, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
        add_textbox(sl, ko,
                    sx, pipe_y + Cm(0.1), stage_w, Cm(0.65),
                    font_size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_textbox(sl, en,
                    sx, pipe_y + Cm(0.7), stage_w, Cm(0.4),
                    font_size=9, italic=True, color=NAVY, align=PP_ALIGN.CENTER)
        # 본문
        add_textbox(sl, sub,
                    sx + Cm(0.3), pipe_y + Cm(1.3), stage_w - Cm(0.6), Cm(0.9),
                    font_size=12, bold=True, color=color, align=PP_ALIGN.LEFT)
        # 데이터 상태
        add_rect(sl, sx + Cm(0.3), pipe_y + pipe_h - Cm(0.8),
                 stage_w - Cm(0.6), Cm(0.5), color)
        add_textbox(sl, "→ " + state,
                    sx + Cm(0.3), pipe_y + pipe_h - Cm(0.8),
                    stage_w - Cm(0.6), Cm(0.5),
                    font_size=10, bold=True, italic=True, color=NAVY, align=PP_ALIGN.CENTER)
        # 화살표
        if i < 3:
            add_chevron(sl, sx + stage_w + Cm(0.0), pipe_y + pipe_h/2 - Cm(0.3),
                        Cm(0.2), Cm(0.7), color)

    # ── 하단: 운영 비용 철학 ──────────────────────────────────────────────
    section2_y = pipe_y + pipe_h + Cm(0.5)
    add_rect(sl, Cm(1.2), section2_y, Cm(0.4), Cm(0.55), ORANGE)
    add_textbox(sl,
                "PART 2.  운영 — 구축의 5배는 운영에서 나간다. 우리는 그 5배를 줄인다.",
                Cm(1.8), section2_y, SLIDE_W - Cm(3), Cm(0.55),
                font_size=12, bold=True, color=ORANGE, align=PP_ALIGN.LEFT)

    # 운영비 폭증 차단 3개 구조
    ops_y = section2_y + Cm(0.7)
    ops_h = Cm(3.6)
    ops_items = [
        ("01", "OpenSearch 기반 자산화",
         "Elastic 상용 라이선스 락인 제거",
         "−40%", "라이선스 비용", CYAN),
        ("02", "Hot · Warm · Cold 티어링",
         "조회 빈도 기반 자동 이관 — 5년 TCO 압축",
         "−60%", "스토리지 TCO", GREEN),
        ("03", "KT Cloud + 자체 모델",
         "외부 토큰 비용 의존 제거 — 운영 예측 가능",
         "3:7 → 7:3", "구축 : 운영", GOLD),
    ]
    ow = (diag_w - 2 * Cm(0.3)) / 3
    for i, (idx, title, body, num, unit, color) in enumerate(ops_items):
        ox = Cm(1.2) + i * (ow + Cm(0.3))
        add_rect(sl, ox, ops_y, ow, ops_h, NAVY2, color, Pt(1.0))
        add_rect(sl, ox, ops_y, Cm(0.18), ops_h, color)

        # 번호 + 제목
        add_textbox(sl, idx,
                    ox + Cm(0.4), ops_y + Cm(0.25), Cm(1.5), Cm(0.6),
                    font_size=11, bold=True, color=color, align=PP_ALIGN.LEFT)
        add_textbox(sl, title,
                    ox + Cm(0.4), ops_y + Cm(0.7), ow - Cm(0.7), Cm(0.8),
                    font_size=15, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        add_textbox(sl, body,
                    ox + Cm(0.4), ops_y + Cm(1.5), ow - Cm(0.7), Cm(1.0),
                    font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)
        # 큰 숫자
        num_y = ops_y + ops_h - Cm(1.4)
        add_rect(sl, ox + Cm(0.3), num_y, ow - Cm(0.6), Cm(1.15), NAVY4, color, Pt(0.5))
        add_textbox(sl, num,
                    ox + Cm(0.3), num_y + Cm(0.05), ow - Cm(0.6), Cm(0.7),
                    font_size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_textbox(sl, unit,
                    ox + Cm(0.3), num_y + Cm(0.7), ow - Cm(0.6), Cm(0.4),
                    font_size=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # ── 하단 결론 띠 ──────────────────────────────────────────────────────
    concl_y = ops_y + ops_h + Cm(0.4)
    add_rect(sl, Cm(1.2), concl_y, SLIDE_W - Cm(2.4), Cm(1.55), GOLD)
    add_textbox(sl,
                "통제 + 지속성 = 진짜 AI 운영 체계",
                Cm(1.5), concl_y + Cm(0.15), SLIDE_W - Cm(3), Cm(0.9),
                font_size=22, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    add_textbox(sl,
                "AI 도입이 아니라, 지속 가능한 AI 운영 체계를 짓습니다  ─  하나은행이 요구한 것은 바로 그것입니다",
                Cm(1.5), concl_y + Cm(0.95), SLIDE_W - Cm(3), Cm(0.55),
                font_size=12, color=NAVY, align=PP_ALIGN.LEFT)

    add_footer(sl, 2)


# ─── 슬라이드 3: 클로징 — 감동 장표 ──────────────────────────────────────────────

def slide_03_closing(prs):
    sl = blank_slide(prs)
    DEEP_NAVY = RGBColor(0x04, 0x09, 0x16)
    fill_bg(sl, DEEP_NAVY)

    cx = SLIDE_W / 2
    cy = SLIDE_H / 2

    # ── 동심원 방사 레이어 (바깥 → 안으로 점점 밝아지는 펄스) ────────────
    pulse = [
        (0.97, 0.97, RGBColor(0x06, 0x0E, 0x20)),
        (0.82, 0.82, RGBColor(0x08, 0x12, 0x28)),
        (0.67, 0.67, RGBColor(0x0B, 0x18, 0x34)),
        (0.52, 0.52, RGBColor(0x0E, 0x1F, 0x42)),
        (0.38, 0.38, RGBColor(0x12, 0x28, 0x50)),
        (0.26, 0.26, RGBColor(0x16, 0x30, 0x5A)),
        (0.15, 0.15, RGBColor(0x1A, 0x38, 0x64)),
    ]
    for rw, rh, color in pulse:
        w = SLIDE_W * rw
        h = SLIDE_H * rh
        sh = sl.shapes.add_shape(MSO_SHAPE.OVAL, cx - w / 2, cy - h / 2, w, h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = color
        sh.line.fill.background()

    # ── 황금빛 궤도 링 (얇은 선 only) ──────────────────────────────────
    for rw, rh, color, lw_pt in [
        (0.72, 0.72, RGBColor(0x30, 0x24, 0x0A), Pt(0.4)),
        (0.52, 0.52, GOLD,                         Pt(0.6)),
        (0.32, 0.32, RGBColor(0x00, 0x6A, 0x80),   Pt(0.3)),
    ]:
        w = SLIDE_W * rw
        h = SLIDE_H * rh
        sh = sl.shapes.add_shape(MSO_SHAPE.OVAL, cx - w / 2, cy - h / 2, w, h)
        sh.fill.background()
        sh.line.color.rgb = color
        sh.line.width = lw_pt

    # ── 중앙 코어 글로우 (GOLD 미세 타원) ───────────────────────────────
    gw, gh = Cm(14), Cm(7)
    sh = sl.shapes.add_shape(MSO_SHAPE.OVAL,
                              cx - gw / 2, cy - gh / 2 - Cm(0.8), gw, gh)
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0x1C, 0x16, 0x04)
    sh.line.fill.background()

    # ── 소형 태그 라인 ───────────────────────────────────────────────────
    add_textbox(sl, "HANA BANK  ×  AI FOUNDATION  ×  KT",
                Cm(0), Cm(4.9), SLIDE_W, Cm(0.55),
                font_size=9, color=RGBColor(0x28, 0x3E, 0x5C), align=PP_ALIGN.CENTER)

    # ── 메인 카피 — 라인 1 (흰색, 보통 굵기) ────────────────────────────
    add_textbox(sl, "통제 가능한 AI 위에서",
                Cm(0), Cm(5.65), SLIDE_W, Cm(1.3),
                font_size=30, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    # ── 메인 카피 — 라인 2 (GOLD, 대형 볼드) ────────────────────────────
    add_textbox(sl, "하나은행의 모든 판단이 움직인다",
                Cm(0), Cm(7.1), SLIDE_W, Cm(1.7),
                font_size=40, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # ── 중앙 구분 장식선 ─────────────────────────────────────────────────
    lw = Cm(5)
    add_rect(sl, cx - lw / 2, Cm(9.1), lw, Cm(0.025), GOLD)
    # 좌우 닷
    d = Cm(0.18)
    add_rect(sl, cx - lw / 2 - d - Cm(0.1), Cm(9.1) - d / 2 + Cm(0.01),
             d, d, GOLD, shape=MSO_SHAPE.OVAL)
    add_rect(sl, cx + lw / 2 + Cm(0.1), Cm(9.1) - d / 2 + Cm(0.01),
             d, d, GOLD, shape=MSO_SHAPE.OVAL)

    # ── 3대 효과 (3열 수평 배치) ─────────────────────────────────────────
    sub_items = [
        (CYAN,   "고객 경험은",  "더 빨라지고"),
        (YELLOW, "직원 업무는",  "더 효율화되며"),
        (ORANGE, "리스크는",    "더 정교하게 통제된다"),
    ]
    col_w = SLIDE_W / 3
    sub_y = Cm(9.65)
    for i, (color, line1, line2) in enumerate(sub_items):
        col_x = col_w * i
        dot = Cm(0.45)
        add_rect(sl, col_x + col_w / 2 - dot / 2, sub_y,
                 dot, dot, color, shape=MSO_SHAPE.OVAL)
        add_textbox(sl, line1,
                    col_x, sub_y + Cm(0.6), col_w, Cm(0.65),
                    font_size=15, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
        add_textbox(sl, line2,
                    col_x, sub_y + Cm(1.2), col_w, Cm(0.65),
                    font_size=15, bold=True, color=color, align=PP_ALIGN.CENTER)

    # ── 수평 thin 구분선 ─────────────────────────────────────────────────
    add_rect(sl, Cm(4), Cm(11.5), SLIDE_W - Cm(8), Cm(0.02),
             RGBColor(0x18, 0x2A, 0x44))

    # ── 변환 카피 (하단, 드라마틱) ──────────────────────────────────────
    tr_y = Cm(12.4)
    # 배경 직사각
    add_rect(sl, 0, tr_y, SLIDE_W, Cm(2.5), RGBColor(0x06, 0x0E, 0x1E))

    # 왼쪽: 도입 (흐린)
    add_textbox(sl, "AI를 도입하는 은행에서",
                Cm(1.0), tr_y + Cm(0.45), SLIDE_W / 2 - Cm(1.5), Cm(0.95),
                font_size=20, color=DARK_TEXT, align=PP_ALIGN.RIGHT)

    # 화살표
    add_textbox(sl, "→",
                cx - Cm(1.0), tr_y + Cm(0.45), Cm(2.0), Cm(0.95),
                font_size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # 오른쪽: 운영 (밝고 볼드)
    add_textbox(sl, "AI를 운영하는 은행으로",
                SLIDE_W / 2 + Cm(0.7), tr_y + Cm(0.45), SLIDE_W / 2 - Cm(1.7), Cm(0.95),
                font_size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # 슬라이드 하단 미니 서명
    add_textbox(sl, "KT B2B 수주전략팀  ·  2026.05",
                Cm(0), SLIDE_H - Cm(0.9), SLIDE_W, Cm(0.65),
                font_size=9, color=RGBColor(0x20, 0x32, 0x4A), align=PP_ALIGN.CENTER)


# ─── 메인 ────────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_war(prs)
    slide_02_operation(prs)
    slide_03_closing(prs)

    out = "/home/user/strategy-pipeline/하나은행_AI_통제하는금융.pptx"
    prs.save(out)
    print(f"✓ Saved: {out}")


if __name__ == "__main__":
    main()
