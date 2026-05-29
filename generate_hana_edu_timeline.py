#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
하나은행 비정형데이터 플랫폼 교육 체계 PPT — 시간 흐름 × 대상자 × 카테고리
출력: /tmp/하나은행_교육체계_타임라인.pptx
"""

from pptx import Presentation
from pptx.util import Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─── 색상 ────────────────────────────────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
SLIDE_BG    = RGBColor(0xFF, 0xFF, 0xFF)
TITLE_COLOR = RGBColor(0x21, 0x3A, 0x6B)
TEXT_COLOR  = RGBColor(0x1F, 0x29, 0x33)
LIGHT_GRAY  = RGBColor(0xF5, 0xF6, 0xF8)
MID_GRAY    = RGBColor(0xE7, 0xEB, 0xF0)
DARK_GRAY   = RGBColor(0x64, 0x74, 0x8B)

PHASE_AMBER = RGBColor(0xF5, 0x9E, 0x0B)   # Phase 배지 배경
PHASE_AMBER_DARK = RGBColor(0xB4, 0x5A, 0x09)

HDR_OPER    = RGBColor(0x1F, 0x4E, 0x9C)   # 플랫폼/시스템 운영자
HDR_ENG     = RGBColor(0x0F, 0x6E, 0x5C)   # 데이터 엔지니어
HDR_BIZ     = RGBColor(0x6B, 0x17, 0x6B)   # 비즈니스 유저

COL_BIFORM  = RGBColor(0x1A, 0x56, 0x8C)   # 비정형 카테고리 배지
COL_OS      = RGBColor(0x0F, 0x6E, 0x5C)   # Object Storage 배지

CELL_BIFORM_BG = RGBColor(0xEB, 0xF2, 0xFA)  # 비정형 셀 배경
CELL_OS_BG     = RGBColor(0xE8, 0xF5, 0xF1)  # OS 셀 배경

FONT_KO = "맑은 고딕"
SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)


# ─── 교육 데이터 ─────────────────────────────────────────────────────────────
PHASES = [
    {
        "num": "01",
        "name": "착수 단계",
        "oper": [
            {
                "cat": "비정형",
                "course": "프로젝트 일반",
                "items": [
                    "데이터 플랫폼 솔루션 교육",
                    "통합 데이터 플랫폼의 구조 이해",
                    "DocuSee 솔루션 플랫폼의 구조 이해",
                    "VectorGreen 솔루션 구조 이해",
                ],
            },
            {
                "cat": "비정형",
                "course": "플랫폼 운영교육",
                "items": [
                    "데이터 수집부터 이메일 및 배포 최적화 가이드",
                    "DocuSee 솔루션 API 사용 방법",
                    "문서 검색 활용을 위한 파일 구조 및 내용",
                    "플랫폼 개요, 사용자 기능",
                ],
            },
            {
                "cat": "OS",
                "course": "Object Storage 솔루션 일반",
                "items": [
                    "MinIO 솔루션 개요 및 구성 이해",
                    "구축 범위 및 운영 절차 설명",
                ],
            },
        ],
        "eng": [],
        "biz": [],
    },
    {
        "num": "02",
        "name": "시스템 설정 채비",
        "oper": [
            {
                "cat": "비정형",
                "course": "플랫폼 운영교육",
                "items": [
                    "VectorGreen 솔루션 UI/API 사용 방법",
                ],
            },
        ],
        "eng": [],
        "biz": [],
    },
    {
        "num": "03",
        "name": "구축 중후반 이후",
        "oper": [
            {
                "cat": "비정형",
                "course": "기술 심화 교육",
                "items": [
                    "수집 인덱싱 배포 파이프라인 운영 및 자동화",
                    "자동화 시나리오 설계 및 모듈 대응",
                    "검색 데이터 구성 및 인덱싱 최적화 개요",
                ],
            },
        ],
        "eng": [
            {
                "cat": "OS",
                "course": "비식별화 기술 교육",
                "items": [
                    "개발환경 이해",
                    "비식별 조치 구성 및 API 활용",
                ],
            },
        ],
        "biz": [],
    },
    {
        "num": "04",
        "name": "구축 완료 후",
        "oper": [
            {
                "cat": "OS",
                "course": "Object Storage 운영 및 장애 대응",
                "items": [
                    "MinIO Cluster 운영 및 상태 확인",
                    "사용자/권한 및 Bucket 관리",
                    "장애 발생 시 대응 절차",
                    "로그 확인 및 기본 점검 방법",
                    "TLS/SSL 및 접근 권한 관리",
                    "Audit Log 및 보안 정책 운영",
                ],
            },
        ],
        "eng": [
            {
                "cat": "OS",
                "course": "Object Storage 기술 교육",
                "items": [
                    "S3 API 연계 및 활용 방법",
                    "Object Storage 구조 및 데이터 관리",
                ],
            },
        ],
        "biz": [],
    },
    {
        "num": "05",
        "name": "운영 전 / 운영 직전",
        "oper": [
            {
                "cat": "OS",
                "course": "비식별화 솔루션 일반",
                "items": [
                    "비식별화 솔루션 운영을 위한 관리자 교육",
                ],
            },
        ],
        "eng": [],
        "biz": [
            {
                "cat": "OS",
                "course": "Object Storage 사용자 교육 (운영 전)",
                "items": [
                    "파일 업로드/다운로드 사용 방법",
                    "기본 기능 및 사용 유의사항 안내",
                ],
            },
            {
                "cat": "OS",
                "course": "비식별화 사용자 교육 (운영 직전)",
                "items": [
                    "개인정보 비식별화 가이드라인",
                    "비식별화 솔루션 사용법 및 기능 실습",
                ],
            },
        ],
    },
    {
        "num": "06",
        "name": "운영 단계",
        "oper": [
            {
                "cat": "비정형",
                "course": "플랫폼 운영교육",
                "items": [
                    "사용자/권한 관리, 시스템 설정",
                    "리소스/워크플로우 모니터링",
                    "정애 대응 및 로그 활용",
                    "사인세 변환최적작을 위한 문서규의 작성가이드",
                    "시스템 아키텍처, 시스템 설정, 모니터링",
                    "로그분석, 장애 코드 식별, 응급조치 방안",
                ],
            },
        ],
        "eng": [
            {
                "cat": "비정형",
                "course": "플랫폼 사용법 (문서 파싱)",
                "items": [
                    "문서 파싱/정형 기능 개요 및 실습",
                    "변환 경로 비교 및 검증 방법",
                    "데이터 변환 관리 개요",
                ],
            },
            {
                "cat": "비정형",
                "course": "플랫폼 사용법 (비정형 수집)",
                "items": [
                    "비정형데이터 수집 및 인덱싱 기능 개요 및 실습",
                    "인덱싱 평가 비교 및 검증 방법",
                    "배포 및 평가 방법",
                ],
            },
        ],
        "biz": [
            {
                "cat": "비정형",
                "course": "플랫폼 사용",
                "items": [
                    "기능 개요, AI 유리 및 보안, 기능 실습",
                    "서비스 분석, 데이터 요청 및 필터링 활용",
                ],
            },
        ],
    },
    {
        "num": "07",
        "name": "안정화 단계",
        "oper": [
            {
                "cat": "비정형",
                "course": "API G/W 운영 교육",
                "items": [
                    "API G/W 기술구조, 설치형성, 상태, 성능 모니터링",
                    "로그 위치 및 포맷, 분석 방법",
                    "장애대응 방법, 백업 및 복구, 증설 방법",
                ],
            },
            {
                "cat": "비정형",
                "course": "Open Search 운영교육",
                "items": [
                    "Open Search 기술 설치, 상태, 성능 모니터링",
                    "모니터링 로그 위치 및 포맷, 분석 방법",
                    "장애대응 방법, 백업 및 복구, 증설 방법",
                ],
            },
        ],
        "eng": [],
        "biz": [
            {
                "cat": "비정형",
                "course": "포털 사용",
                "items": [
                    "권한별 메뉴 구조 및 주요 기능",
                ],
            },
        ],
    },
]


# ─── 헬퍼 ────────────────────────────────────────────────────────────────────

def blank_slide(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = SLIDE_BG
    return sl


def add_rect(slide, x, y, w, h, bg_color, text="", font_size=10,
             bold=False, font_color=None, align=PP_ALIGN.LEFT,
             line_color=None, valign=None, wrap=True, italic=False):
    shape = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Emu(9525)
    else:
        shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = wrap
        if valign:
            tf.vertical_anchor = valign
        tf.margin_left  = Cm(0.18)
        tf.margin_right = Cm(0.18)
        tf.margin_top   = Cm(0.08)
        tf.margin_bottom = Cm(0.08)
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = FONT_KO
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = font_color if font_color else TEXT_COLOR

    return shape


def add_multiline(slide, x, y, w, h, bg_color, lines, font_size=8.5,
                  line_color=None, cat_color=None):
    """과정명 + 교육내용 bullet을 하나의 텍스트박스에 출력."""
    from pptx.util import Pt as PPt
    shape = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Emu(9525)
    else:
        shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left   = Cm(0.18)
    tf.margin_right  = Cm(0.18)
    tf.margin_top    = Cm(0.1)
    tf.margin_bottom = Cm(0.1)

    for i, (text, is_course) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.name = FONT_KO
        run.font.size = PPt(font_size if not is_course else font_size + 0.5)
        run.font.bold = is_course
        run.font.color.rgb = (cat_color if cat_color and is_course else TEXT_COLOR)
    return shape


def add_cat_badge(slide, x, y, cat):
    color = COL_BIFORM if cat == "비정형" else COL_OS
    label = "비정형" if cat == "비정형" else "Object Storage"
    add_rect(slide, x, y, 2.8, 0.42, color,
             text=label, font_size=7.5, bold=True,
             font_color=WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


# ─── 슬라이드 1: 타이틀 ───────────────────────────────────────────────────────

def make_title_slide(prs):
    sl = blank_slide(prs)
    # 배경 상단 띠
    add_rect(sl, 0, 0, 33.87, 4.5, TITLE_COLOR)
    # 하나은행 로고 영역
    add_rect(sl, 1.2, 0.6, 6, 0.7, TITLE_COLOR,
             text="하나은행", font_size=20, bold=True, font_color=WHITE)
    # 메인 타이틀
    add_rect(sl, 1.2, 1.5, 31, 1.5, TITLE_COLOR,
             text="비정형데이터 플랫폼 구축사업 교육 체계",
             font_size=26, bold=True, font_color=WHITE, align=PP_ALIGN.LEFT)
    add_rect(sl, 1.2, 3.0, 31, 0.9, TITLE_COLOR,
             text="시간 흐름 × 교육 대상자 × 카테고리별 분류",
             font_size=15, bold=False, font_color=RGBColor(0xB8, 0xC8, 0xE0),
             align=PP_ALIGN.LEFT)

    # 범례 박스
    legend_y = 5.5
    add_rect(sl, 1.2, legend_y, 30, 0.5, LIGHT_GRAY,
             text="교육 카테고리 범례", font_size=10, bold=True,
             font_color=TITLE_COLOR)
    add_rect(sl, 1.2, legend_y + 0.55, 4.5, 0.6, COL_BIFORM,
             text="■  비정형 데이터 처리", font_size=9.5, bold=True,
             font_color=WHITE, valign=MSO_ANCHOR.MIDDLE)
    add_rect(sl, 6.0, legend_y + 0.55, 4.5, 0.6, COL_OS,
             text="■  Object Storage", font_size=9.5, bold=True,
             font_color=WHITE, valign=MSO_ANCHOR.MIDDLE)

    # 교육 대상자 컬럼 색상 범례
    add_rect(sl, 1.2, legend_y + 1.4, 4.5, 0.55, HDR_OPER,
             text="플랫폼/시스템 운영자", font_size=9, bold=True,
             font_color=WHITE, valign=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    add_rect(sl, 6.0, legend_y + 1.4, 4.5, 0.55, HDR_ENG,
             text="데이터 엔지니어", font_size=9, bold=True,
             font_color=WHITE, valign=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    add_rect(sl, 10.8, legend_y + 1.4, 4.5, 0.55, HDR_BIZ,
             text="비즈니스 유저", font_size=9, bold=True,
             font_color=WHITE, valign=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    # Phase 타임라인 요약
    phases_short = ["01\n착수 단계", "02\n시스템\n설정 채비", "03\n구축\n중후반 이후",
                    "04\n구축\n완료 후", "05\n운영 전\n운영 직전", "06\n운영 단계", "07\n안정화\n단계"]
    tl_y = legend_y + 2.5
    add_rect(sl, 1.2, tl_y, 30, 0.4, MID_GRAY,
             text="교육 시간 흐름 (Phase 01 → Phase 07)", font_size=9.5, bold=True,
             font_color=TITLE_COLOR, valign=MSO_ANCHOR.MIDDLE)
    box_w = 30 / 7
    for i, ph in enumerate(phases_short):
        color = PHASE_AMBER if i % 2 == 0 else RGBColor(0xE0, 0x85, 0x07)
        add_rect(sl, 1.2 + i * box_w, tl_y + 0.45, box_w - 0.08, 1.5, color,
                 text=ph, font_size=8, bold=True,
                 font_color=WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # 슬라이드 총 수 안내
    add_rect(sl, 1.2, tl_y + 2.1, 30, 0.5, LIGHT_GRAY,
             text="본 자료는 총 9장으로 구성됩니다  (타이틀 1장 + 로드맵 개요 1장 + Phase별 상세 7장)",
             font_size=9, bold=False, font_color=DARK_GRAY,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


# ─── 슬라이드 2: 전체 로드맵 개요 ────────────────────────────────────────────

def make_overview_slide(prs):
    sl = blank_slide(prs)
    # 타이틀 바
    add_rect(sl, 0, 0, 33.87, 1.1, TITLE_COLOR,
             text="전체 교육 로드맵 개요  |  비정형데이터 플랫폼 구축사업",
             font_size=13, bold=True, font_color=WHITE, valign=MSO_ANCHOR.MIDDLE)

    # 컬럼 헤더
    col_headers = ["Phase", "교육 시기", "플랫폼/시스템 운영자", "데이터 엔지니어", "비즈니스 유저"]
    col_x  = [0.2,  2.5,   8.0,   17.5,   26.0]
    col_w  = [2.1,  5.3,   9.3,    8.3,    7.7]
    col_bg = [PHASE_AMBER, MID_GRAY, HDR_OPER, HDR_ENG, HDR_BIZ]
    hdr_fc = [WHITE, TITLE_COLOR, WHITE, WHITE, WHITE]

    for i, (hdr, cx, cw, bg, fc) in enumerate(zip(col_headers, col_x, col_w, col_bg, hdr_fc)):
        add_rect(sl, cx, 1.15, cw, 0.65, bg,
                 text=hdr, font_size=9.5, bold=True,
                 font_color=fc, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
                 line_color=WHITE)

    # 데이터 행
    row_h = (19.05 - 1.15 - 0.65 - 0.2) / 7
    phase_names = [p["name"] for p in PHASES]

    def summarize(entries):
        if not entries:
            return "—"
        parts = []
        for e in entries:
            badge = "[비정형]" if e["cat"] == "비정형" else "[OS]"
            parts.append(f"{badge} {e['course']}")
        return "\n".join(parts)

    for r, ph in enumerate(PHASES):
        ry = 1.15 + 0.65 + r * row_h
        row_bg = LIGHT_GRAY if r % 2 == 0 else WHITE

        # Phase 번호
        add_rect(sl, 0.2, ry, 2.1, row_h, PHASE_AMBER,
                 text=f"PHASE\n{ph['num']}", font_size=9, bold=True,
                 font_color=WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
                 line_color=WHITE)
        # 시기
        add_rect(sl, 2.5, ry, 5.3, row_h, row_bg,
                 text=ph["name"], font_size=9, bold=True,
                 font_color=TITLE_COLOR, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
                 line_color=MID_GRAY)
        # 운영자
        add_rect(sl, 8.0, ry, 9.3, row_h, CELL_BIFORM_BG if any(e["cat"] == "비정형" for e in ph["oper"]) else CELL_OS_BG if ph["oper"] else row_bg,
                 text=summarize(ph["oper"]), font_size=8, bold=False,
                 font_color=TEXT_COLOR, line_color=MID_GRAY, valign=MSO_ANCHOR.MIDDLE)
        # 엔지니어
        add_rect(sl, 17.5, ry, 8.3, row_h, CELL_BIFORM_BG if any(e["cat"] == "비정형" for e in ph["eng"]) else CELL_OS_BG if ph["eng"] else row_bg,
                 text=summarize(ph["eng"]), font_size=8, bold=False,
                 font_color=TEXT_COLOR, line_color=MID_GRAY, valign=MSO_ANCHOR.MIDDLE)
        # 비즈니스 유저
        add_rect(sl, 26.0, ry, 7.7, row_h, CELL_BIFORM_BG if any(e["cat"] == "비정형" for e in ph["biz"]) else CELL_OS_BG if ph["biz"] else row_bg,
                 text=summarize(ph["biz"]), font_size=8, bold=False,
                 font_color=TEXT_COLOR, line_color=MID_GRAY, valign=MSO_ANCHOR.MIDDLE)


# ─── 슬라이드 3~9: Phase별 상세 ──────────────────────────────────────────────

def make_phase_slide(prs, ph):
    sl = blank_slide(prs)

    # 헤더
    add_rect(sl, 0, 0, 33.87, 1.0, TITLE_COLOR,
             text=f"PHASE {ph['num']}  |  {ph['name']}  —  교육 대상자별 상세 커리큘럼",
             font_size=13, bold=True, font_color=WHITE, valign=MSO_ANCHOR.MIDDLE)

    # Phase 진행 위치 표시
    box_w_ph = 33.87 / 7
    for i, p in enumerate(PHASES):
        is_cur = p["num"] == ph["num"]
        color  = PHASE_AMBER if is_cur else MID_GRAY
        fc     = WHITE if is_cur else DARK_GRAY
        add_rect(sl, i * box_w_ph, 1.05, box_w_ph - 0.04, 0.55, color,
                 text=f"P{p['num']} {p['name']}", font_size=7, bold=is_cur,
                 font_color=fc, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
                 line_color=WHITE)

    # 컬럼 헤더
    CONTENT_TOP = 1.65
    COL_SPECS = [
        ("플랫폼 / 시스템 운영자", HDR_OPER, "oper"),
        ("데이터 엔지니어",         HDR_ENG,  "eng"),
        ("비즈니스 유저",           HDR_BIZ,  "biz"),
    ]
    col_w = (33.87 - 0.4) / 3
    hdr_h = 0.65

    for ci, (label, bg, key) in enumerate(COL_SPECS):
        cx = 0.2 + ci * col_w
        add_rect(sl, cx, CONTENT_TOP, col_w - 0.1, hdr_h, bg,
                 text=label, font_size=11, bold=True,
                 font_color=WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
                 line_color=WHITE)

    # 컨텐츠 영역
    BODY_TOP = CONTENT_TOP + hdr_h + 0.1
    BODY_H   = 19.05 - BODY_TOP - 0.15

    for ci, (label, bg, key) in enumerate(COL_SPECS):
        cx = 0.2 + ci * col_w
        entries = ph[key]

        if not entries:
            add_rect(sl, cx, BODY_TOP, col_w - 0.1, BODY_H, LIGHT_GRAY,
                     text="해당 Phase에 배정된 교육 없음",
                     font_size=9, bold=False, font_color=DARK_GRAY,
                     align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
                     line_color=MID_GRAY)
            continue

        # 각 과정 블록을 세로로 나눔
        block_h = BODY_H / len(entries)
        for ei, entry in enumerate(entries):
            ey = BODY_TOP + ei * block_h
            cat = entry["cat"]
            cell_bg = CELL_BIFORM_BG if cat == "비정형" else CELL_OS_BG
            cat_color = COL_BIFORM if cat == "비정형" else COL_OS

            # 카테고리 배지
            badge_h = 0.38
            add_rect(sl, cx, ey, 3.0, badge_h,
                     cat_color,
                     text=("비정형 데이터 처리" if cat == "비정형" else "Object Storage"),
                     font_size=7.5, bold=True, font_color=WHITE,
                     align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

            # 과정명 셀
            course_h = 0.5
            add_rect(sl, cx, ey + badge_h, col_w - 0.1, course_h, MID_GRAY,
                     text=entry["course"], font_size=9.5, bold=True,
                     font_color=TITLE_COLOR, valign=MSO_ANCHOR.MIDDLE,
                     line_color=MID_GRAY)

            # 교육 내용 셀
            items_h = block_h - badge_h - course_h
            items_text = "\n".join(f"• {it}" for it in entry["items"])
            add_rect(sl, cx, ey + badge_h + course_h, col_w - 0.1, items_h,
                     cell_bg, text=items_text, font_size=8.5, bold=False,
                     font_color=TEXT_COLOR, valign=MSO_ANCHOR.TOP,
                     line_color=MID_GRAY)


# ─── 메인 ────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    make_title_slide(prs)
    make_overview_slide(prs)
    for ph in PHASES:
        make_phase_slide(prs, ph)

    out = "/tmp/하나은행_교육체계_타임라인.pptx"
    prs.save(out)
    print(f"저장 완료: {out}  ({len(prs.slides)}장)")


if __name__ == "__main__":
    main()
