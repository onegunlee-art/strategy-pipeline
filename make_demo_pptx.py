from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Color palette
C_BG      = RGBColor(0x0A, 0x0E, 0x1A)   # 딥 네이비
C_CARD    = RGBColor(0x12, 0x1A, 0x2E)   # 카드 배경
C_ACCENT  = RGBColor(0x00, 0xD4, 0xFF)   # 시안 강조
C_GREEN   = RGBColor(0x00, 0xFF, 0xA3)   # 그린
C_YELLOW  = RGBColor(0xFF, 0xD6, 0x00)   # 옐로
C_RED     = RGBColor(0xFF, 0x4D, 0x6D)   # 레드
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY    = RGBColor(0x8A, 0x9B, 0xB8)
C_BORDER  = RGBColor(0x1E, 0x2D, 0x4A)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank

def add_slide():
    return prs.slides.add_slide(blank_layout)

def bg(slide, color=C_BG):
    from pptx.util import Emu
    bg = slide.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def box(slide, l, t, w, h, fill=C_CARD, alpha=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = C_BORDER
    s.line.width = Pt(0.5)
    return s

def txt(slide, text, l, t, w, h, size=18, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Malgun Gothic"
    return tb

def accent_line(slide, l, t, w, color=C_ACCENT, thickness=3):
    line = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Pt(thickness))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

def chip(slide, label, l, t, w=1.4, h=0.35, color=C_ACCENT):
    s = slide.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))  # rounded rect
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    tf = s.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = label
    run.font.size = Pt(11)
    run.font.color.rgb = C_BG
    run.font.bold = True
    run.font.name = "Malgun Gothic"

def multiline_txt(slide, lines, l, t, w, h, size=16, color=C_WHITE, bold=False, line_spacing=1.2):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = "Malgun Gothic"

# ─────────────────────────────────────────────
# SLIDE 1: 표지
# ─────────────────────────────────────────────
sl = add_slide()
bg(sl)

# 대형 배경 그라데이션 원
circle = sl.shapes.add_shape(9, Inches(7.5), Inches(-1), Inches(8), Inches(8))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x00, 0x3A, 0x5C)
circle.line.fill.background()

accent_line(sl, 1.0, 3.8, 4.0, C_ACCENT)

txt(sl, "전략 인텔리전스 데모", 1.0, 1.5, 9, 1.2, size=42, bold=True, color=C_WHITE)
txt(sl, "수주전략  ×  지정학 분석", 1.0, 2.8, 9, 0.8, size=26, color=C_ACCENT)
txt(sl, '"같은 엔진으로 두 가지 불확실성을 다룬다"', 1.0, 4.2, 10, 0.7, size=18, color=C_GRAY)
txt(sl, "2026.06", 1.0, 6.6, 3, 0.5, size=14, color=C_GRAY)

# ─────────────────────────────────────────────
# SLIDE 2: 문제 정의
# ─────────────────────────────────────────────
sl = add_slide()
bg(sl)
accent_line(sl, 0.8, 0.6, 2.5, C_RED)
txt(sl, "우리가 해결하는 질문", 0.8, 0.8, 11, 0.8, size=30, bold=True)

questions = [
    ('① ', '"이 딜, 지금 몇 % 가능성인가?"'),
    ('② ', '"지정학 이슈가 우리 사업에 미치는 영향은?"'),
    ('③ ', '"임원들이 같은 방향을 보고 있는가?"'),
]
for i, (num, q) in enumerate(questions):
    y = 2.0 + i * 1.3
    box(sl, 0.8, y, 11.6, 1.0, fill=C_CARD)
    txt(sl, num, 1.1, y + 0.15, 0.5, 0.7, size=20, color=C_ACCENT, bold=True)
    txt(sl, q, 1.6, y + 0.15, 10.5, 0.7, size=20, color=C_WHITE)

box(sl, 0.8, 5.8, 11.6, 0.8, fill=RGBColor(0x0A, 0x2A, 0x1A))
txt(sl, "→  데이터는 있다.  판단 구조가 없다.", 1.2, 5.95, 10, 0.5, size=20, color=C_GREEN, bold=True)

# ─────────────────────────────────────────────
# SLIDE 3: 시스템 구조
# ─────────────────────────────────────────────
sl = add_slide()
bg(sl)
accent_line(sl, 0.8, 0.6, 2.5, C_ACCENT)
txt(sl, "시스템 구조", 0.8, 0.8, 11, 0.7, size=30, bold=True)

# GIST box
box(sl, 0.5, 1.8, 4.5, 4.5, fill=RGBColor(0x0A, 0x1A, 0x2A))
txt(sl, "the gist.", 0.8, 2.0, 4.0, 0.5, size=20, bold=True, color=C_ACCENT)
multiline_txt(sl, ["• 300+ 기사 RAG", "• 충돌점 / 일치점 분석", "• AI 인사이트 자동생성", "• 고객 서비스 중 — 읽기만"], 0.8, 2.7, 3.8, 2.5, size=16, color=C_GRAY)
txt(sl, "결과만 전송 (read-only)", 0.8, 5.5, 4.0, 0.5, size=13, color=C_GREEN)

# Arrow
txt(sl, "→", 5.2, 3.5, 0.8, 0.6, size=36, color=C_ACCENT, align=PP_ALIGN.CENTER)
txt(sl, "API", 5.1, 4.1, 1.0, 0.4, size=12, color=C_GRAY, align=PP_ALIGN.CENTER)

# Dashboard box
box(sl, 6.3, 1.8, 6.5, 4.5, fill=RGBColor(0x0A, 0x2A, 0x1A))
txt(sl, "전략 대시보드", 6.6, 2.0, 5.8, 0.5, size=20, bold=True, color=C_GREEN)
multiline_txt(sl, ["• 드라이버 점수 추출", "• 확률 분포 계산", "• 임원 시그널 버튼", "• Polymarket 비교", "• 리포트 생성"], 6.6, 2.7, 5.8, 2.5, size=16, color=C_GRAY)
txt(sl, "수주전략 동일 엔진 재사용", 6.6, 5.5, 5.8, 0.5, size=13, color=C_YELLOW)

# ─────────────────────────────────────────────
# SLIDE 4: 데모 플로우
# ─────────────────────────────────────────────
sl = add_slide()
bg(sl)
accent_line(sl, 0.8, 0.6, 2.5, C_YELLOW)
txt(sl, "데모 플로우 (12분)", 0.8, 0.8, 11, 0.7, size=30, bold=True)

steps = [
    (C_ACCENT,  "① 주제 입력",      '"이란 전쟼 종전 가능성"'),
    (C_GRAY,    "② GIST 분석",      "기사 20건 · 충돌점 6 · 일치점 14"),
    (C_GREEN,   "③ AI 초안",        "드라이버 점수 · 확률 22% · 범위 ±12%p"),
    (C_YELLOW,  "④ 비교",           "Polymarket 34%  vs  우리 22%"),
    (C_ACCENT,  "⑤ 블라인드 투표",  "QR 스캔 · 시그널 버튼 클릭"),
    (C_RED,     "⑥ 공개",           "분포 갱신 · 이견 가시화  ← WOW"),
    (C_GREEN,   "⑦ 인사이트",       "가장 중요한 드라이버 표시"),
    (C_YELLOW,  "⑧ 딜 연결",        "우리 수주 파이프라인 영향  ← WOW"),
    (C_WHITE,   "⑨ 리포트 생성",    "SG 제주 포맷 · PDF 다운로드"),
]

for i, (color, step, desc) in enumerate(steps):
    col = i // 5
    row = i % 5
    x = 0.5 + col * 6.6
    y = 1.8 + row * 1.08
    box(sl, x, y, 6.3, 0.9, fill=C_CARD)
    txt(sl, step, x + 0.15, y + 0.12, 2.2, 0.65, size=15, color=color, bold=True)
    txt(sl, desc, x + 2.4, y + 0.15, 3.7, 0.6, size=14, color=C_GRAY)

# ─────────────────────────────────────────────
# SLIDE 5: 드라이버 + 확률 분포
# ─────────────────────────────────────────────
sl = add_slide()
bg(sl)
accent_line(sl, 0.8, 0.6, 3.0, C_GREEN)
txt(sl, "화면 ①  드라이버 + 확률 분포", 0.8, 0.8, 11, 0.7, size=28, bold=True)

# Driver panel
box(sl, 0.5, 1.8, 6.0, 5.0, fill=C_CARD)
txt(sl, "드라이버 분석 (AI 초안)", 0.8, 2.0, 5.5, 0.5, size=16, bold=True, color=C_ACCENT)

drivers = [
    ("외교 채널",  3, C_RED),
    ("군사 강도",  7, C_RED),
    ("경제 압박",  6, C_YELLOW),
    ("이란 내부",  3, C_RED),
    ("호르무즈",   4, C_YELLOW),
]
for i, (name, score, color) in enumerate(drivers):
    y = 2.7 + i * 0.72
    txt(sl, name, 0.8, y, 2.0, 0.5, size=14, color=C_WHITE)
    # bar background
    bar_bg = sl.shapes.add_shape(1, Inches(2.9), Inches(y + 0.08), Inches(2.8), Inches(0.28))
    bar_bg.fill.solid(); bar_bg.fill.fore_color.rgb = RGBColor(0x1E, 0x2D, 0x4A)
    bar_bg.line.fill.background()
    # bar fill
    bar_w = 2.8 * score / 10
    bar = sl.shapes.add_shape(1, Inches(2.9), Inches(y + 0.08), Inches(bar_w), Inches(0.28))
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    txt(sl, f"{score}/10", 5.85, y, 0.6, 0.4, size=13, color=C_GRAY)

txt(sl, "충돌 기사 6건 → 불확실성 ±12%p 자동 반영", 0.7, 6.3, 5.8, 0.4, size=12, color=C_GRAY)

# Probability distribution panel
box(sl, 7.0, 1.8, 5.8, 5.0, fill=C_CARD)
txt(sl, "종전 가능성", 7.3, 2.0, 5.2, 0.5, size=16, bold=True, color=C_ACCENT)

# Draw simple ASCII-style distribution bars
bar_data = [1, 2, 4, 7, 9, 8, 6, 4, 2, 1]
bar_w_each = 0.38
for i, h in enumerate(bar_data):
    bh = h * 0.28
    bx = 7.4 + i * bar_w_each
    by = 5.5 - bh
    b = sl.shapes.add_shape(1, Inches(bx), Inches(by), Inches(bar_w_each - 0.04), Inches(bh))
    b.fill.solid()
    b.fill.fore_color.rgb = C_ACCENT if i in [3,4,5,6] else RGBColor(0x1E, 0x3A, 0x5C)
    b.line.fill.background()

txt(sl, "8%", 7.3, 5.6, 0.6, 0.4, size=12, color=C_GRAY)
txt(sl, "22%", 8.8, 5.6, 0.7, 0.4, size=14, color=C_ACCENT, bold=True)
txt(sl, "34%", 11.3, 5.6, 0.7, 0.4, size=12, color=C_GRAY)
box(sl, 7.0, 6.1, 5.8, 0.55, fill=RGBColor(0x0A, 0x2A, 0x1A))
txt(sl, "범위: ±12%p  |  충돌 기사 많을수록 넓어짐", 7.2, 6.2, 5.4, 0.4, size=13, color=C_GREEN)

# ─────────────────────────────────────────────
# SLIDE 6: Polymarket 비교
# ─────────────────────────────────────────────
sl = add_slide()
bg(sl)
accent_line(sl, 0.8, 0.6, 3.0, C_YELLOW)
txt(sl, "화면 ②  Polymarket 비교", 0.8, 0.8, 11, 0.7, size=28, bold=True)

# Our analysis
box(sl, 0.5, 1.8, 5.8, 3.5, fill=C_CARD)
txt(sl, "우리 분석", 0.8, 2.0, 5.2, 0.5, size=16, bold=True, color=C_ACCENT)
txt(sl, "22%", 1.5, 2.6, 4.0, 1.4, size=72, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
txt(sl, "±12%p", 2.0, 4.0, 2.8, 0.5, size=18, color=C_GRAY, align=PP_ALIGN.CENTER)

# VS
txt(sl, "vs", 6.4, 3.0, 0.8, 0.8, size=28, color=C_GRAY, bold=True, align=PP_ALIGN.CENTER)

# Polymarket
box(sl, 7.3, 1.8, 5.5, 3.5, fill=C_CARD)
txt(sl, "Polymarket", 7.6, 2.0, 5.0, 0.5, size=16, bold=True, color=C_YELLOW)
txt(sl, "34%", 7.8, 2.6, 4.0, 1.4, size=72, bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)
txt(sl, "시장 집단지성", 8.3, 4.0, 2.8, 0.5, size=16, color=C_GRAY, align=PP_ALIGN.CENTER)

# Insight box
box(sl, 0.5, 5.6, 12.3, 1.4, fill=RGBColor(0x0A, 0x2A, 0x1A))
txt(sl, "차이  -12%p  →  호르무즈 레버리지 변수를 시장이 아직 미반영", 0.9, 5.75, 11.5, 0.5, size=17, color=C_GREEN, bold=True)
txt(sl, "이것이 우리 분석의 엣지 (Judgement Moat)", 0.9, 6.2, 11.5, 0.5, size=15, color=C_GRAY)

# ─────────────────────────────────────────────
# SLIDE 7: 블라인드 투표
# ─────────────────────────────────────────────
sl = add_slide()
bg(sl)
accent_line(sl, 0.8, 0.6, 3.0, C_ACCENT)
txt(sl, "화면 ③  블라인드 투표", 0.8, 0.8, 11, 0.7, size=28, bold=True)

# Projector screen simulation
box(sl, 0.5, 1.7, 6.2, 5.2, fill=C_CARD)
txt(sl, "프로젝터 화면", 0.8, 1.85, 5.5, 0.4, size=13, color=C_GRAY)
# QR placeholder
qr_box = sl.shapes.add_shape(1, Inches(2.0), Inches(2.4), Inches(3.0), Inches(2.5))
qr_box.fill.solid(); qr_box.fill.fore_color.rgb = RGBColor(0x1E, 0x2D, 0x4A)
qr_box.line.color.rgb = C_ACCENT; qr_box.line.width = Pt(1.5)
txt(sl, "QR", 3.0, 3.3, 1.0, 0.8, size=32, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
txt(sl, "각자 폰으로 스캔하세요", 1.0, 5.1, 5.0, 0.4, size=13, color=C_GRAY, align=PP_ALIGN.CENTER)
txt(sl, "참여: 3/5명  ⏱ 01:20", 1.0, 5.55, 5.0, 0.4, size=14, color=C_YELLOW, align=PP_ALIGN.CENTER)
txt(sl, "결과는 모두 제출 후 공개", 1.0, 5.95, 5.0, 0.35, size=12, color=C_GRAY, align=PP_ALIGN.CENTER)

# Mobile screen simulation
box(sl, 7.2, 1.7, 5.6, 5.2, fill=C_CARD)
txt(sl, "임원 폰 화면", 7.5, 1.85, 5.0, 0.4, size=13, color=C_GRAY)
txt(sl, "이란 전쟼 종전 가능성", 7.5, 2.3, 5.0, 0.45, size=15, bold=True, color=C_WHITE)
txt(sl, "AI 추정: 22%", 7.5, 2.75, 5.0, 0.35, size=13, color=C_ACCENT)

buttons = [
    ("협상 재개 소식", C_GREEN),
    ("협상 결렬 선언", C_RED),
    ("공습 중단 발표", C_GREEN),
    ("공습 빈도 증가", C_RED),
    ("제재 완화 신호", C_GREEN),
    ("호르무즈 봉쇄 강화", C_RED),
]
for i, (label, color) in enumerate(buttons):
    col = i % 2
    row = i // 2
    bx = 7.4 + col * 2.65
    by = 3.25 + row * 0.75
    btn = sl.shapes.add_shape(5, Inches(bx), Inches(by), Inches(2.45), Inches(0.55))
    btn.fill.solid(); btn.fill.fore_color.rgb = RGBColor(0x1E, 0x2D, 0x4A)
    btn.line.color.rgb = color; btn.line.width = Pt(1.5)
    tf = btn.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = label
    run.font.size = Pt(12)
    run.font.color.rgb = color
    run.font.name = "Malgun Gothic"

# ─────────────────────────────────────────────
# SLIDE 8: 결과 공개 — WOW
# ─────────────────────────────────────────────
sl = add_slide()
bg(sl)
accent_line(sl, 0.8, 0.6, 4.0, C_RED)
txt(sl, "화면 ④  결과 공개  ★ 핵심 WOW", 0.8, 0.8, 11, 0.7, size=28, bold=True, color=C_WHITE)

# Vote results
box(sl, 0.5, 1.7, 5.8, 4.8, fill=C_CARD)
txt(sl, "투표 결과", 0.8, 1.9, 5.2, 0.45, size=16, bold=True, color=C_ACCENT)

votes = [
    ("김 부사장", "협상 재개 소식", C_GREEN, "▲"),
    ("이 본부장", "공습 빈도 증가", C_RED,   "▼"),
    ("박 상무  ", "제재 완화 신호", C_GREEN, "▲"),
    ("최 팀장  ", "강경파 집권 신호", C_RED, "▼"),
    ("정 상무  ", "협상 재개 소식", C_GREEN, "▲"),
]
for i, (name, signal, color, arrow) in enumerate(votes):
    y = 2.5 + i * 0.72
    txt(sl, name, 0.8, y, 1.8, 0.55, size=14, color=C_GRAY)
    txt(sl, signal, 2.6, y, 2.8, 0.55, size=14, color=color)
    txt(sl, arrow, 5.5, y, 0.5, 0.55, size=16, color=color, bold=True)

box(sl, 0.5, 6.05, 5.8, 0.5, fill=RGBColor(0x0A, 0x2A, 0x1A))
txt(sl, "낙관 3명 / 비관 2명  →  이견 자체가 정보", 0.8, 6.12, 5.2, 0.4, size=13, color=C_GREEN)

# Distribution change
box(sl, 6.8, 1.7, 6.0, 4.8, fill=C_CARD)
txt(sl, "확률 분포 변화", 7.1, 1.9, 5.4, 0.45, size=16, bold=True, color=C_ACCENT)

# Before
txt(sl, "투표 전", 7.1, 2.5, 2.0, 0.4, size=13, color=C_GRAY)
bar_data_before = [1, 2, 4, 7, 9, 8, 6, 4, 2, 1]
for i, h in enumerate(bar_data_before):
    bh = h * 0.2
    bx = 7.1 + i * 0.3
    by = 4.5 - bh
    b = sl.shapes.add_shape(1, Inches(bx), Inches(by), Inches(0.26), Inches(bh))
    b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5C)
    b.line.fill.background()
txt(sl, "22%  ±12%p", 7.1, 4.55, 2.5, 0.4, size=12, color=C_GRAY)

# Arrow
txt(sl, "→", 10.1, 3.5, 0.6, 0.6, size=28, color=C_ACCENT, bold=True)

# After
txt(sl, "투표 후", 10.9, 2.5, 2.0, 0.4, size=13, color=C_GREEN)
bar_data_after = [1, 2, 3, 6, 12, 11, 7, 3, 1, 1]
for i, h in enumerate(bar_data_after):
    bh = h * 0.18
    bx = 10.9 + i * 0.28
    by = 4.5 - bh
    b = sl.shapes.add_shape(1, Inches(bx), Inches(by), Inches(0.24), Inches(bh))
    b.fill.solid(); b.fill.fore_color.rgb = C_GREEN if i in [3,4,5,6] else RGBColor(0x1E, 0x3A, 0x5C)
    b.line.fill.background()
txt(sl, "26%  ±7%p", 10.9, 4.55, 2.5, 0.4, size=12, color=C_GREEN)

box(sl, 6.8, 6.05, 6.0, 0.5, fill=RGBColor(0x0A, 0x2A, 0x1A))
txt(sl, "22% → 26%  |  불확실성 ±12%p → ±7%p 감소", 7.0, 6.12, 5.6, 0.4, size=13, color=C_GREEN)

# ─────────────────────────────────────────────
# SLIDE 9: 가장 중요한 드라이버
# ─────────────────────────────────────────────
sl = add_slide()
bg(sl)
accent_line(sl, 0.8, 0.6, 3.0, C_GREEN)
txt(sl, "화면 ⑤  지금 무엇을 봐야 하는가", 0.8, 0.8, 11, 0.7, size=28, bold=True)

sensitivity = [
    ("외교 채널",  "+1점",  "+9%p",  "★★★", "최우선 모니터링", C_RED),
    ("군사 강도",  "-1점",  "+6%p",  "★★",  "공습 현황 추적",  C_YELLOW),
    ("호르무즈",   "-1점",  "+5%p",  "★★",  "해협 통항 데이터", C_YELLOW),
    ("경제 압박",  "+1점",  "+2%p",  "★",   "제재 동향 참고",  C_GRAY),
]

for i, (driver, change, delta, stars, action, color) in enumerate(sensitivity):
    y = 1.9 + i * 1.28
    box(sl, 0.5, y, 12.3, 1.1, fill=C_CARD)
    txt(sl, driver,  0.9,  y+0.25, 2.2, 0.6, size=17, bold=True, color=C_WHITE)
    txt(sl, change,  3.2,  y+0.25, 1.3, 0.6, size=16, color=color)
    txt(sl, "→",     4.5,  y+0.25, 0.5, 0.6, size=16, color=C_GRAY)
    txt(sl, delta,   5.0,  y+0.25, 1.4, 0.6, size=20, bold=True, color=color)
    txt(sl, stars,   6.5,  y+0.25, 1.5, 0.6, size=16, color=C_YELLOW)
    txt(sl, action,  8.2,  y+0.25, 4.3, 0.6, size=15, color=C_GRAY)

box(sl, 0.5, 7.0, 12.3, 0.5, fill=RGBColor(0x00, 0x2A, 0x1A))
txt(sl, "→  협상 테이블 뉴스 하나가 확률을 9%p 움직입니다", 0.9, 7.08, 11.5, 0.38, size=15, color=C_GREEN, bold=True)

# ─────────────────────────────────────────────
# SLIDE 10: 딜 파이프라인 연결
# ─────────────────────────────────────────────
sl = add_slide()
bg(sl)
accent_line(sl, 0.8, 0.6, 4.0, C_YELLOW)
txt(sl, "화면 ⑥  지정학 → 수주전략 직결  ★ WOW", 0.8, 0.8, 11, 0.7, size=28, bold=True)

box(sl, 0.5, 1.7, 5.5, 1.0, fill=RGBColor(0x0A, 0x2A, 0x3A))
txt(sl, "이란 전쟼 종전 가능성  26%  반영 시", 0.8, 1.85, 5.0, 0.65, size=16, color=C_ACCENT, bold=True)

deals_impact = [
    ("제주 위성데이터센터", "에너지 공급 리스크", "▲",  C_RED),
    ("xx 물류인프라 사업",  "납기 리스크",        "▲ +8%p", C_RED),
    ("yy 해저케이블 사업",  "착공 지연 리스크",   "▲",  C_YELLOW),
]

for i, (deal, risk, change, color) in enumerate(deals_impact):
    y = 3.0 + i * 1.3
    box(sl, 0.5, y, 12.3, 1.1, fill=C_CARD)
    txt(sl, deal,   0.9, y+0.25, 4.0, 0.6, size=17, bold=True, color=C_WHITE)
    txt(sl, risk,   5.0, y+0.25, 4.5, 0.6, size=16, color=C_GRAY)
    txt(sl, change, 9.6, y+0.25, 2.5, 0.6, size=18, bold=True, color=color)

box(sl, 0.5, 6.8, 12.3, 0.6, fill=RGBColor(0x1A, 0x2A, 0x0A))
txt(sl, '"지정학 분석이 수주전략의 입력값이 된다"', 0.9, 6.9, 11.5, 0.42, size=16, color=C_GREEN, bold=True)

# ─────────────────────────────────────────────
# SLIDE 11: 왜 이게 다른가
# ─────────────────────────────────────────────
sl = add_slide()
bg(sl)
accent_line(sl, 0.8, 0.6, 3.0, C_ACCENT)
txt(sl, "왜 이게 다른가", 0.8, 0.8, 11, 0.7, size=30, bold=True)

# Headers
box(sl, 0.5, 1.7, 5.8, 0.6, fill=RGBColor(0x1E, 0x1A, 0x1A))
box(sl, 6.8, 1.7, 5.8, 0.6, fill=RGBColor(0x0A, 0x2A, 0x1A))
txt(sl, "기존 방식", 1.8, 1.8, 3.5, 0.45, size=16, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
txt(sl, "이 시스템", 8.1, 1.8, 3.5, 0.45, size=16, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

comparisons = [
    ("뉴스 보고 느낌으로 판단",    "드라이버로 구조화된 판단"),
    ("개인 직관",                  "집단 판단 + 이견 가시화"),
    ("단일 숫자",                  "분포 + 신뢰 구간"),
    ("Polymarket 모름",            "Polymarket 실시간 비교"),
    ("지정학 따로, 수주 따로",      "지정학 → 수주 자동 연결"),
    ("리포트 없음",                "즉시 발간 · PDF · 이메일"),
]

for i, (before, after) in enumerate(comparisons):
    y = 2.5 + i * 0.77
    box(sl, 0.5, y, 5.8, 0.65, fill=C_CARD)
    box(sl, 6.8, y, 5.8, 0.65, fill=C_CARD)
    txt(sl, "✗  " + before, 0.8, y+0.12, 5.2, 0.45, size=14, color=C_GRAY)
    txt(sl, "✓  " + after,  7.1, y+0.12, 5.2, 0.45, size=14, color=C_GREEN)

# ─────────────────────────────────────────────
# SLIDE 12: 다음 단계
# ─────────────────────────────────────────────
sl = add_slide()
bg(sl)
accent_line(sl, 0.8, 0.6, 2.5, C_ACCENT)
txt(sl, "다음 단계", 0.8, 0.8, 11, 0.7, size=30, bold=True)

phases = [
    ("Phase 1", "지금 가능",   "이란 전쟼 데모 완성 · SG 제주 포맷 리포트",  C_ACCENT),
    ("Phase 2", "+2주",        "GIST API 연결 · 실데이터 파이프라인",         C_GREEN),
    ("Phase 3", "+1개월",      "임원 투표 → 수주전략 자동 반영",             C_YELLOW),
    ("Phase 4", "+3개월",      "전딜 × 전이슈 실시간 워룸",                  C_RED),
]

for i, (phase, timing, desc, color) in enumerate(phases):
    y = 1.8 + i * 1.3
    box(sl, 0.5, y, 12.3, 1.1, fill=C_CARD)
    chip_s = sl.shapes.add_shape(5, Inches(0.7), Inches(y+0.25), Inches(1.5), Inches(0.6))
    chip_s.fill.solid(); chip_s.fill.fore_color.rgb = color
    chip_s.line.fill.background()
    tf = chip_s.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run(); run.text = phase
    run.font.size = Pt(13); run.font.color.rgb = C_BG
    run.font.bold = True; run.font.name = "Malgun Gothic"

    txt(sl, timing, 2.5,  y+0.25, 1.5, 0.6, size=15, color=C_GRAY)
    txt(sl, desc,   4.2,  y+0.25, 8.2, 0.6, size=15, color=C_WHITE)

box(sl, 0.5, 7.0, 12.3, 0.5, fill=RGBColor(0x00, 0x1A, 0x2A))
txt(sl, '"오늘 데모는 Phase 1 입니다"', 1.0, 7.08, 11.3, 0.38, size=15, color=C_ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────
out_path = "/home/user/strategy-pipeline/전략인텔리전스_데모_발표자료.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
