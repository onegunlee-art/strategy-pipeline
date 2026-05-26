#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
하나은행 비정형데이터 자산화 플랫폼 수주전략 PPT 생성
위성영상_ISP_수주전략.pptx 레이아웃·색상·폰트 그대로 복제
"""

from pptx import Presentation
from pptx.util import Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── 색상 (위성영상 PPT에서 추출) ────────────────────────────────────────────
HEADER_BG   = RGBColor(0xE7, 0xEB, 0xF0)
CELL_BG     = RGBColor(0xFF, 0xFF, 0xFF)
TITLE_COLOR = RGBColor(0x21, 0x3A, 0x6B)
TEXT_COLOR  = RGBColor(0x1F, 0x29, 0x33)
RED_BADGE   = RGBColor(0xC0, 0x1B, 0x2E)
BLUE_BADGE  = RGBColor(0x1F, 0x4E, 0x9C)
GREEN_BADGE = RGBColor(0x17, 0x6B, 0x2E)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
SLIDE_BG    = RGBColor(0xFF, 0xFF, 0xFF)

FONT_KO = "맑은 고딕"
SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)

# ─── 공통 헬퍼 ───────────────────────────────────────────────────────────────

def blank_slide(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = SLIDE_BG
    return sl


def add_rect(slide, x, y, w, h, bg_color, text="", font_size=10.5,
             bold=False, font_color=None, align=PP_ALIGN.LEFT,
             line_color=None, valign=None):
    from pptx.enum.text import MSO_ANCHOR
    shape = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Emu(9525)
    else:
        shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        if valign:
            tf.vertical_anchor = valign
        # 내부 여백
        tf.margin_left = Cm(0.2)
        tf.margin_right = Cm(0.2)
        tf.margin_top = Cm(0.1)
        tf.margin_bottom = Cm(0.1)

        lines = text.split("\n")
        for i, line in enumerate(lines):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            para.alignment = align
            run = para.add_run()
            run.text = line
            run.font.name = FONT_KO
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = font_color or TEXT_COLOR
    return shape


def add_textbox(slide, x, y, w, h, text, font_size=12, bold=False,
                font_color=None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.name = FONT_KO
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = font_color or TEXT_COLOR
    return txBox


def add_badge(slide, x, y, text, color):
    add_rect(slide, x, y, 1.8, 0.65, color,
             text=text, font_size=9, bold=True, font_color=WHITE,
             align=PP_ALIGN.CENTER)


def add_slide_title(slide, text):
    add_textbox(slide, 1.27, 0.5, 30.48, 1.6, text,
                font_size=24, bold=True, font_color=TITLE_COLOR)


# ─── 슬라이드 1: 과업범위 ────────────────────────────────────────────────────

def slide1_scope(prs):
    sl = blank_slide(prs)

    # 좌측 테이블 헤더
    add_rect(sl, 1.27, 2.8, 3.56, 1.0, HEADER_BG,
             "구  분", 10.5, bold=True)
    add_rect(sl, 4.83, 2.8, 16.51, 1.0, HEADER_BG,
             "사업내용", 10.5, bold=True)

    rows = [
        ("데이터\n인프라",
         "○ 데이터 수집 파이프라인 구축 (원천 연계, 출처태그 자동 부착)\n"
         "○ Hot·Warm·Cold 티어 스토리지 아키텍처 구축\n"
         "○ 금융문서 특화 파싱·청킹·임베딩 처리 파이프라인 (HWP·PDF·Oracle CLOB 파서 내장)"),
        ("플랫폼\n개발",
         "○ 워크플로우 엔진 및 품질관리 시스템 개발\n"
         "○ 메타데이터·리니지·버전 관리 체계 구축\n"
         "○ 검색 포털·API/SDK·API Gateway 개발\n"
         "○ 정형/비정형 실시간 비식별화 처리 (펜타시스템 협력)"),
        ("AI 연계",
         "○ MCP Server 기반 멀티에이전트 연동 아키텍처 구현\n"
         "○ 하이브리드 시맨틱-벡터 검색 엔진 구축 (BD 협력)\n"
         "○ 영업점 상담 Q&A·여신심사 비정형 요약·AML 탐지 등 킬러앱 즉시 가동"),
    ]

    heights = [3.3, 4.0, 4.0]
    y = 3.8
    for (label, content), h in zip(rows, heights):
        add_rect(sl, 1.27, y, 3.56, h, CELL_BG, label, 10.5, bold=True)
        add_rect(sl, 4.83, y, 16.51, h, CELL_BG, content, 9.5)
        y += h

    # 우측 추진방향
    add_textbox(sl, 22.35, 2.0, 10.67, 1.4, "추진 방향",
                font_size=22, bold=True, font_color=TITLE_COLOR)

    bullets = [
        "• 파편화된 비정형 데이터를 AI가 즉시 활용 가능한\n'AI-Ready All-Data' 체계로 전환",
        "• Hot·Warm·Cold 티어링으로\n5년 스토리지 TCO 60% 절감",
        "• 개인정보보호법·신용정보법·전금법 완벽 준수\n— 망분리·DRM·비식별화 내재화",
        "• MCP 기반 에이전트 아키텍처로 구축 즉시\n현업이 자율적 AI 서비스 확장 가능",
    ]
    by = 4.0
    for b in bullets:
        add_textbox(sl, 22.35, by, 11.18, 2.6, b, font_size=11,
                    font_color=TEXT_COLOR)
        by += 3.3

    return sl


# ─── 슬라이드 2: 고객/사업 현황 분석 ─────────────────────────────────────────

def slide2_customer(prs):
    sl = blank_slide(prs)
    add_slide_title(sl, "1. 고객/사업 현황 분석")

    # 헤더
    add_rect(sl, 1.27, 3.1, 5.59, 1.0, HEADER_BG, "과업범위", 10.5, bold=True)
    add_rect(sl, 6.86, 3.1, 25.74, 1.0, HEADER_BG, "고객니즈/페인포인트", 10.5, bold=True)

    rows = [
        ("데이터\n인프라",
         "• [데이터 파편화] 약관·상품설명서·여신서류 등 비정형 데이터 사일로 — AI 학습 불가 구조 직면\n"
         "• [DRM 복호화] 내부 DRM 걸린 문서 복호화 처리 및 감사로그 필요\n"
         "• [망분리 환경] 개발망·운영망 분리 환경에서 온프레미스 임베딩 모델 운용 필수"),
        ("플랫폼\n개발",
         "• [레거시 연계] Oracle CLOB/BLOB 등 기존 코어시스템과의 안정적 데이터 연계 요구\n"
         "• [컴플라이언스] 개인정보보호법·신용정보법·전금법 대응 비식별화 체계 구축 필요\n"
         "• [데이터 이관] 초기 1TB/100만 건 + 일 10GB/4,000건 증분 수용 가능한 확장형 아키텍처"),
        ("AI 연계",
         "• [AI 확장성] 구축 후 현업이 자율적으로 AI 서비스 확장 가능한 오픈 환경 요구\n"
         "• [업무 자동화] 영업점 Q&A·여신심사 요약·AML 탐지 등 실무 킬러앱 즉시 적용 원함\n"
         "• [운영 안정성] 장기 운영비(구축비의 5배) 절감 및 지속 가능한 운영체계 요구"),
    ]

    heights = [3.6, 3.6, 3.6]
    y = 4.1
    for (label, content), h in zip(rows, heights):
        add_rect(sl, 1.27, y, 5.59, h, CELL_BG, label, 10.5, bold=True)
        add_rect(sl, 6.86, y, 25.74, h, CELL_BG, content, 9.5)
        y += h

    # 하단 footer
    add_textbox(sl, 1.27, 17.5, 30.48, 0.9,
                "* 평가 위원 현황: 하나은행 IT·디지털혁신·데이터·컴플라이언스 담당 부서 실무자 구성 예상",
                font_size=9, font_color=TEXT_COLOR)

    return sl


# ─── 슬라이드 3: 경쟁 현황 분석 ──────────────────────────────────────────────

def slide3_competition(prs):
    sl = blank_slide(prs)
    add_slide_title(sl, "2. 경쟁 현황 분석")

    col1_x, col1_w = 1.27, 5.08
    col2_x, col2_w = 6.35, 13.0
    col3_x, col3_w = 19.35, 13.0

    # 헤더
    add_rect(sl, col1_x, 3.1, col1_w, 1.0, HEADER_BG, "과업범위", 10.5, bold=True)
    add_rect(sl, col2_x, 3.1, col2_w, 1.0, HEADER_BG, "KT", 12, bold=True)
    add_rect(sl, col3_x, 3.1, col3_w, 1.0, HEADER_BG, "삼성SDS", 12, bold=True)

    rows = [
        ("데이터·\n플랫폼\n개발",
         "컨소시엄 수행 예상 영역\n협력: weda(파이프라인)·BD(검색)·펜타시스템(비식별화)\n"
         "• ES 기반 글로벌 금융 표준 — 락인 리스크 없음\n"
         "• HWP·CLOB 파서, DRM 복호화 API 내재화\n"
         "• 단, ES 라이선스 비용으로 입찰가 여유 제한적",
         "주관사 수행 예상 영역\nBrightics AI 자체솔루션 기반\n"
         "• 라이선스 비용 없음 → 입찰가 낮출 가능성 ↑\n"
         "• 금융 특화 파서·비식별화 별도 조달 필요\n"
         "• 자체솔루션 유지보수 종속 리스크",
         "열위(가격)", "우위(가격)"),
        ("인프라·\n운영",
         "주관사 수행 예상 영역\n"
         "• KT 통신 인프라 운영 DNA 기반 99.9% SLA\n"
         "• Hot/Warm/Cold 티어링 TCO 60% 절감 비교표\n"
         "• 온프레미스 임베딩 → 외부 토큰 비용 제거\n"
         "• 망분리 환경 안정적 운영 레퍼런스 보유",
         "컨소시엄 수행 예상 영역\n"
         "• 대형 금융 인프라 레퍼런스 보유\n"
         "• 자사솔루션 기반 운영 종속 구조\n"
         "• 외부 LLM 의존 시 운영비 예측 어려움",
         "우위(운영)", "열위(운영)"),
        ("AI\n연계",
         "• MCP Server 기반 즉시 가동 — 킬러앱 청사진 제시\n"
         "• 구축 즉시 영업점 Q&A·여신심사·AML 적용\n"
         "• 규제 변경 → 내규 매핑 자동화 선제 제안\n"
         "• 장기 선진 로드맵(NGA 모델) 컨설팅 제공",
         "• 플랫폼 구축 중심, AI 비즈니스 청사진 부재 가능성\n"
         "• Brightics 에코시스템 내 에이전트 확장 제한\n"
         "• ISP 수행 시 기간연장·비용 증가 이슈 선례",
         "우위(확장성)", "열위(확장성)"),
    ]

    heights = [4.3, 4.3, 4.0]
    y = 4.1
    for row, h in zip(rows, heights):
        label, kt_txt, sds_txt, kt_badge, sds_badge = row
        add_rect(sl, col1_x, y, col1_w, h, CELL_BG, label, 10.5, bold=True)
        add_rect(sl, col2_x, y, col2_w, h, CELL_BG, kt_txt, 8.5)
        add_rect(sl, col3_x, y, col3_w, h, CELL_BG, sds_txt, 8.5)
        # badge: KT
        badge_color_kt = BLUE_BADGE if "우위" in kt_badge else RED_BADGE
        add_badge(sl, col2_x + col2_w - 2.0, y + 0.15, kt_badge, badge_color_kt)
        # badge: SDS
        badge_color_sds = BLUE_BADGE if "우위" in sds_badge else RED_BADGE
        add_badge(sl, col3_x + col3_w - 2.0, y + 0.15, sds_badge, badge_color_sds)
        y += h

    # LG CNS 비고
    add_textbox(sl, 1.27, 17.2, 30.48, 0.9,
                "* LG CNS: 참여 의지 약함 (영업 정보) — 사실상 KT vs 삼성SDS 2파전 구도",
                font_size=9, font_color=TEXT_COLOR)

    return sl


# ─── 슬라이드 4: 차별화 아이템 도출 ──────────────────────────────────────────

def slide4_diff(prs):
    sl = blank_slide(prs)
    add_slide_title(sl, "3. 차별화 아이템 도출")

    add_rect(sl, 1.27, 3.1, 5.59, 1.0, HEADER_BG, "과업범위", 10.5, bold=True)
    add_rect(sl, 6.86, 3.1, 25.74, 1.0, HEADER_BG, "KT 차별화 아이템", 10.5, bold=True)

    rows = [
        ("데이터·\n플랫폼\n개발",
         "• [오픈소스 표준] ES 기반 글로벌 금융 표준(JPMorgan·Goldman Sachs) — 삼성SDS 자체솔루션 대비 락인 리스크 Zero\n"
         "• [금융 특화 파서] HWP 파서 내장·Oracle CLOB/BLOB 직연계·DRM 복호화 API+감사로그 기본 제공\n"
         "• [2-Pass 클렌징] 입수 시점 + AI 처리 후 이중 품질 검증 → 오류율 0.5% 이하 보장"),
        ("인프라·\n운영",
         "• [TCO 절감] Hot/Warm/Cold 티어링 → 스토리지 5년 TCO 60% 절감 비교표 + KT Cloud 할인율\n"
         "   - Hot(SSD/ES): 최근 6개월만 인덱싱 / Warm(HDD): 6개월~2년 / Cold(S3): 2년 이상\n"
         "• [운영 DNA] 통신 인프라 운영 노하우 기반 99.9% SLA — '멈추지 않는 운영'\n"
         "• [비용 예측] KT Cloud + 온프레미스 임베딩 모델 → 외부 LLM 토큰 비용 의존 제거, 구축:운영 비율 개선"),
        ("AI\n연계",
         "• [즉시 가동] MCP Server → 구축 즉시 영업점 Q&A·여신심사 요약·AML 탐지 3대 킬러앱 적용 (추가 제안)\n"
         "• [컴플라이언스 자동화] 규제 변경 감지 → 내규 영향도 자동 매핑 → 감사 대응 선제 제안\n"
         "• [장기 로드맵] 선진 금융 데이터 생태계(NGA 모델) 전환 로드맵 및 위성 전문가 컨설팅 제공"),
    ]

    heights = [3.6, 4.2, 3.8]
    y = 4.1
    for (label, content), h in zip(rows, heights):
        add_rect(sl, 1.27, y, 5.59, h, CELL_BG, label, 10.5, bold=True)
        add_rect(sl, 6.86, y, 25.74, h, CELL_BG, content, 9.5)
        y += h

    return sl


# ─── 메인 ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide1_scope(prs)
    slide2_customer(prs)
    slide3_competition(prs)
    slide4_diff(prs)

    out = "하나은행_비정형데이터_수주전략.pptx"
    prs.save(out)
    print(f"✅ 저장 완료: {out}")
    print(f"   슬라이드 수: {len(prs.slides)}")


if __name__ == "__main__":
    main()
