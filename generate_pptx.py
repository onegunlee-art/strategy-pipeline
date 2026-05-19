#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
금융권 AI Agent 시장 조사 보고서 PPT 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ─── 색상 팔레트 ────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0A, 0x16, 0x28)      # 배경 진한 네이비
NAVY2       = RGBColor(0x0D, 0x1F, 0x3C)      # 약간 밝은 네이비 (카드 BG)
CYAN        = RGBColor(0x00, 0xD4, 0xFF)      # 청록 강조
ORANGE      = RGBColor(0xFF, 0x6B, 0x35)      # 주황 강조
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xD6, 0xE0)
YELLOW      = RGBColor(0xFF, 0xD7, 0x00)
GREEN       = RGBColor(0x00, 0xE0, 0x96)
MID_BLUE    = RGBColor(0x1A, 0x4A, 0x7A)      # 카드 테두리
DARK_TEXT   = RGBColor(0xA0, 0xB8, 0xCC)      # 본문 흐린 텍스트

FONT_KO = "맑은 고딕"

# 슬라이드 크기 16:9
SLIDE_W = Cm(33.87)
SLIDE_H = Cm(19.05)

FOOTER_TEXT = "국내 금융권 AI Agent 실증 분석  |  2026.05  |  CONFIDENTIAL"

# ─── 유틸리티 함수들 ─────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color: RGBColor):
    """슬라이드 배경색 채우기"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color: RGBColor, line_color=None, line_width=None):
    """직사각형 도형 추가"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, font_name=FONT_KO,
                italic=False, word_wrap=True):
    """텍스트 박스 추가 (단일 단락)"""
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = word_wrap
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_footer(slide):
    """하단 풋라인 추가"""
    # 구분선
    add_rect(slide,
             Cm(1.5), SLIDE_H - Cm(1.2),
             SLIDE_W - Cm(3), Cm(0.04),
             RGBColor(0x1A, 0x4A, 0x7A))
    add_textbox(slide, FOOTER_TEXT,
                Cm(1.5), SLIDE_H - Cm(1.1),
                SLIDE_W - Cm(3), Cm(0.9),
                font_size=8, color=DARK_TEXT, align=PP_ALIGN.CENTER)


def add_slide_number(slide, num):
    add_textbox(slide, str(num),
                SLIDE_W - Cm(2.5), SLIDE_H - Cm(1.1),
                Cm(2), Cm(0.9),
                font_size=8, color=DARK_TEXT, align=PP_ALIGN.RIGHT)


def add_section_tag(slide, text, x, y, bg=CYAN, fg=NAVY):
    """섹션 태그 배지"""
    w, h = Cm(5.5), Cm(0.65)
    r = add_rect(slide, x, y, w, h, bg)
    add_textbox(slide, text, x, y, w, h,
                font_size=9, bold=True, color=fg, align=PP_ALIGN.CENTER)
    return w, h


def add_percent_badge(slide, pct_text, x, y, color: RGBColor):
    """성공률 배지"""
    w, h = Cm(3.2), Cm(1.1)
    add_rect(slide, x, y, w, h, color)
    add_textbox(slide, pct_text, x, y, w, h,
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def add_bullet_card(slide, title, bullets, x, y, w, h,
                    title_color=CYAN, bg_color=NAVY2,
                    border_color=MID_BLUE):
    """카드형 박스 (제목 + 불릿 리스트)"""
    add_rect(slide, x, y, w, h, bg_color, border_color, Pt(0.75))
    # 제목
    add_textbox(slide, title,
                x + Cm(0.3), y + Cm(0.2),
                w - Cm(0.6), Cm(0.7),
                font_size=11, bold=True, color=title_color)
    # 불릿
    by = y + Cm(0.9)
    for b in bullets:
        add_textbox(slide, f"▸  {b}",
                    x + Cm(0.4), by,
                    w - Cm(0.7), Cm(0.55),
                    font_size=9, color=LIGHT_GRAY)
        by += Cm(0.52)


def add_kpi_block(slide, label, value, unit, x, y, w=Cm(4.5), h=Cm(2.4),
                  val_color=CYAN):
    """KPI 숫자 블록"""
    add_rect(slide, x, y, w, h, NAVY2, MID_BLUE, Pt(0.75))
    add_textbox(slide, value,
                x, y + Cm(0.3), w, Cm(1.2),
                font_size=28, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_textbox(slide, unit,
                x, y + Cm(1.3), w, Cm(0.55),
                font_size=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide, label,
                x, y + Cm(1.85), w, Cm(0.55),
                font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ─── 슬라이드 1 : 표지 ───────────────────────────────────────────────────────────

def slide_01_cover(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)

    # 좌측 강조 세로선
    add_rect(sl, Cm(1.4), Cm(2.5), Cm(0.3), Cm(10), CYAN)

    # 상단 배지
    add_rect(sl, Cm(2.2), Cm(2.5), Cm(8), Cm(0.7), MID_BLUE)
    add_textbox(sl, "MARKET RESEARCH REPORT  |  AI AGENT",
                Cm(2.2), Cm(2.5), Cm(8), Cm(0.7),
                font_size=9, color=CYAN, align=PP_ALIGN.LEFT)

    # 메인 제목
    add_textbox(sl, "국내 금융권\nAI Agent\n시장 현황 및\n활용 사례",
                Cm(2.2), Cm(3.6), Cm(22), Cm(8),
                font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # 부제목
    add_textbox(sl, "K금융그룹 · S금융그룹 파일럿 기반 실증 분석",
                Cm(2.2), Cm(12.0), Cm(22), Cm(1.0),
                font_size=16, color=CYAN, align=PP_ALIGN.LEFT)

    # 날짜 및 조직
    add_rect(sl, Cm(2.2), Cm(13.4), Cm(10), Cm(0.06), CYAN)
    add_textbox(sl, "2026.05",
                Cm(2.2), Cm(13.6), Cm(6), Cm(0.8),
                font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

    # 우측 장식 – 원형 배경
    add_rect(sl, Cm(23), Cm(1), Cm(10), Cm(17), RGBColor(0x0D, 0x1F, 0x3C),
             MID_BLUE, Pt(0.5))

    # 우측 통계 블록
    add_kpi_block(sl, "조사 금융그룹", "2", "개 그룹", Cm(24.5), Cm(2.5))
    add_kpi_block(sl, "분석 영역", "6", "개 영역", Cm(24.5), Cm(5.3), val_color=ORANGE)
    add_kpi_block(sl, "파일럿 기간", "4", "개월", Cm(24.5), Cm(8.1))
    add_kpi_block(sl, "최고 성공률", "92.9", "%", Cm(24.5), Cm(10.9), val_color=GREEN)

    add_footer(sl)


# ─── 슬라이드 2 : 목차 ───────────────────────────────────────────────────────────

def slide_02_toc(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)

    add_rect(sl, Cm(1.4), Cm(1.2), Cm(0.3), Cm(1.5), CYAN)
    add_textbox(sl, "목  차", Cm(2.0), Cm(1.2), Cm(10), Cm(1.5),
                font_size=28, bold=True, color=WHITE)
    add_textbox(sl, "TABLE OF CONTENTS",
                Cm(2.0), Cm(2.5), Cm(12), Cm(0.6),
                font_size=10, color=CYAN)

    items = [
        ("01", "조사 개요"),
        ("02", "금융 AI Agent 6대 적용 영역"),
        ("03", "영역별 성공률 분석"),
        ("04", "고객/마케팅  —  최고 성공 영역 (92.9%)"),
        ("05", "문서·콘텐츠 생성 (77.3%)"),
        ("06", "상담 지원 (70.0%)"),
        ("07", "투자/리서치 (42.9%)"),
        ("08", "리스크/보안 (37.5%)"),
        ("09", "내부 운영 지원 (33.3%)"),
        ("10", "시사점 및 전략적 제언"),
        ("11", "하나은행 적용 방향"),
    ]

    cols = [items[:6], items[6:]]
    col_x = [Cm(2.0), Cm(18.0)]

    for ci, col in enumerate(cols):
        cx = col_x[ci]
        for i, (num, title) in enumerate(col):
            y = Cm(3.6) + i * Cm(2.2)
            # 번호 배지
            add_rect(sl, cx, y, Cm(1.3), Cm(0.85),
                     CYAN if ci == 0 else ORANGE)
            add_textbox(sl, num, cx, y, Cm(1.3), Cm(0.85),
                        font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
            # 제목
            add_textbox(sl, title,
                        cx + Cm(1.5), y + Cm(0.05),
                        Cm(13), Cm(0.8),
                        font_size=12, color=WHITE)
            # 구분선
            add_rect(sl, cx, y + Cm(0.9), Cm(14.5), Cm(0.02),
                     RGBColor(0x1A, 0x4A, 0x7A))

    add_footer(sl)
    add_slide_number(sl, 2)


# ─── 슬라이드 3 : 조사 개요 ──────────────────────────────────────────────────────

def slide_03_overview(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)

    add_rect(sl, Cm(1.4), Cm(1.2), Cm(0.3), Cm(1.5), CYAN)
    add_textbox(sl, "조사 개요", Cm(2.0), Cm(1.2), Cm(15), Cm(1.5),
                font_size=28, bold=True, color=WHITE)
    add_section_tag(sl, "RESEARCH OVERVIEW", Cm(2.0), Cm(2.6))

    cards = [
        ("조사 대상", CYAN, [
            "K금융그룹 (국내 주요 금융그룹)",
            "S금융그룹 (국내 주요 금융그룹)",
            "은행·증권·보험 계열사 포함",
        ]),
        ("조사 기간", ORANGE, [
            "2024년 ~ 2025년 상반기",
            "4개월 파일럿 프로그램 운영",
            "실제 업무 환경 적용 기준",
        ]),
        ("분석 항목", GREEN, [
            "AI Agent 도입 유형별 성공률",
            "6개 주요 적용 영역 분류",
            "세부 기능 단위 성과 분석",
        ]),
        ("성공 기준", YELLOW, [
            "파일럿 완료 후 본 사업 전환 비율",
            "정량적 KPI 달성 여부",
            "현업 부서 지속 사용 의향 포함",
        ]),
    ]

    for i, (title, color, bullets) in enumerate(cards):
        col = i % 2
        row = i // 2
        cx = Cm(2.0) + col * Cm(16.0)
        cy = Cm(4.0) + row * Cm(6.0)
        cw, ch = Cm(14.5), Cm(5.5)

        add_rect(sl, cx, cy, cw, ch, NAVY2, color, Pt(1.5))
        # 좌측 색상 바
        add_rect(sl, cx, cy, Cm(0.3), ch, color)
        add_textbox(sl, title,
                    cx + Cm(0.7), cy + Cm(0.3),
                    cw - Cm(0.9), Cm(0.9),
                    font_size=14, bold=True, color=color)
        add_rect(sl, cx + Cm(0.7), cy + Cm(1.2),
                 cw - Cm(0.9), Cm(0.03), color)
        for j, b in enumerate(bullets):
            add_textbox(sl, f"•  {b}",
                        cx + Cm(0.8), cy + Cm(1.4) + j * Cm(1.2),
                        cw - Cm(1.0), Cm(1.0),
                        font_size=11, color=LIGHT_GRAY)

    add_footer(sl)
    add_slide_number(sl, 3)


# ─── 슬라이드 4 : 6대 영역 성공률 차트 ──────────────────────────────────────────

def slide_04_chart(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)

    add_rect(sl, Cm(1.4), Cm(1.2), Cm(0.3), Cm(1.5), CYAN)
    add_textbox(sl, "6대 영역 성공률 한눈에", Cm(2.0), Cm(1.2), Cm(20), Cm(1.5),
                font_size=28, bold=True, color=WHITE)
    add_section_tag(sl, "SUCCESS RATE OVERVIEW", Cm(2.0), Cm(2.6))

    # 차트 데이터
    categories = [
        "고객/마케팅",
        "문서·콘텐츠 생성",
        "상담 지원",
        "투자/리서치",
        "리스크/보안",
        "내부 운영 지원",
    ]
    values = [92.9, 77.3, 70.0, 42.9, 37.5, 33.3]
    colors_bar = [
        RGBColor(0x00, 0xD4, 0xFF),  # CYAN  ≥70
        RGBColor(0x00, 0xD4, 0xFF),
        RGBColor(0x00, 0xD4, 0xFF),
        RGBColor(0xFF, 0xD7, 0x00),  # YELLOW 50-70
        RGBColor(0xFF, 0x6B, 0x35),  # ORANGE <50
        RGBColor(0xFF, 0x6B, 0x35),
    ]

    chart_data = ChartData()
    chart_data.categories = categories
    chart_data.add_series("성공률 (%)", values)

    chart_shape = sl.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Cm(1.8), Cm(3.4), Cm(20), Cm(13.5),
        chart_data,
    )
    chart = chart_shape.chart

    # 차트 배경: XML 직접 조작으로 배경색 설정
    cs_elem = chart._element  # c:chartSpace
    # chartSpace 배경을 NAVY2로 설정 (spPr 추가)
    ns_c  = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    ns_a  = "http://schemas.openxmlformats.org/drawingml/2006/main"
    navy2_hex = "0D1F3C"

    def _set_solid_spPr(parent_elem, hex_color):
        """parent_elem 에 <c:spPr> solid fill 추가 (없으면 생성)"""
        spPr = parent_elem.find(qn("c:spPr"))
        if spPr is None:
            spPr = etree.SubElement(parent_elem, qn("c:spPr"))
        # 기존 solidFill 제거
        for child in list(spPr):
            spPr.remove(child)
        solidFill = etree.SubElement(spPr, qn("a:solidFill"))
        srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", hex_color)
        return spPr

    _set_solid_spPr(cs_elem, navy2_hex)

    # plotArea 배경도 NAVY2
    c_chart = cs_elem.find(qn("c:chart"))
    c_plotArea = c_chart.find(qn("c:plotArea"))
    _set_solid_spPr(c_plotArea, navy2_hex)

    plot = chart.plots[0]
    series = plot.series[0]

    # 각 막대 색상 개별 설정
    for idx, pt in enumerate(series.points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colors_bar[idx]
        pt.format.line.fill.background()

    # 데이터 레이블
    series.data_labels.show_value = True
    series.data_labels.number_format = '0.0"%"'
    series.data_labels.font.size = Pt(11)
    series.data_labels.font.bold = True
    series.data_labels.font.color.rgb = WHITE

    # 축 설정
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(11)
    cat_axis.tick_labels.font.color.rgb = WHITE
    cat_axis.tick_labels.font.name = FONT_KO
    cat_axis.format.line.fill.background()

    val_axis = chart.value_axis
    val_axis.maximum_scale = 100
    val_axis.tick_labels.font.size = Pt(9)
    val_axis.tick_labels.font.color.rgb = DARK_TEXT
    val_axis.format.line.fill.background()
    val_axis.major_gridlines.format.line.color.rgb = RGBColor(0x1A, 0x3A, 0x5A)

    chart.has_legend = False

    # 우측 범례 설명
    legend_data = [
        (CYAN,   "≥ 70%   고성공 영역"),
        (YELLOW, "50~69%  중간 영역"),
        (ORANGE, "< 50%   개선 필요"),
    ]
    for i, (c, label) in enumerate(legend_data):
        lx = Cm(23.0)
        ly = Cm(5.0) + i * Cm(1.8)
        add_rect(sl, lx, ly, Cm(0.7), Cm(0.55), c)
        add_textbox(sl, label, lx + Cm(0.9), ly, Cm(9), Cm(0.6),
                    font_size=11, color=LIGHT_GRAY)

    add_footer(sl)
    add_slide_number(sl, 4)


# ─── 슬라이드 5 : 고객/마케팅 ────────────────────────────────────────────────────

def slide_05_marketing(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)

    add_rect(sl, Cm(1.4), Cm(1.2), Cm(0.3), Cm(1.5), CYAN)
    add_textbox(sl, "고객 / 마케팅", Cm(2.0), Cm(1.2), Cm(18), Cm(1.5),
                font_size=28, bold=True, color=WHITE)
    add_section_tag(sl, "CUSTOMER & MARKETING", Cm(2.0), Cm(2.65))

    # 성공률 배지
    add_percent_badge(sl, "92.9%", Cm(25.5), Cm(1.3), CYAN)
    add_textbox(sl, "최고 성공 영역", Cm(25.0), Cm(2.55), Cm(4.5), Cm(0.6),
                font_size=9, color=CYAN, align=PP_ALIGN.CENTER)

    # 왜 최고 성공률인가 — 설명 카드
    reasons = [
        ("즉시 가시적 효과 측정", "클릭률, 반응률 등 정량 KPI 즉시 측정 가능\n마케터 대비 10배 이상 콘텐츠 생산량 달성"),
        ("규제 리스크 최소화", "내부 활용 중심으로 금융 규제 충돌 없음\n외부 공시 전 인간 검토 단계 유지"),
        ("명확한 업무 범위", "상품 광고 카피, FAQ, 이벤트 메시지 등\n구체적이고 반복적인 태스크 중심 적용"),
    ]
    for i, (title, desc) in enumerate(reasons):
        rx = Cm(2.0) + i * Cm(10.5)
        ry = Cm(3.6)
        rw, rh = Cm(9.8), Cm(4.0)
        add_rect(sl, rx, ry, rw, rh, NAVY2, CYAN, Pt(1.0))
        add_rect(sl, rx, ry, rw, Cm(0.12), CYAN)
        add_textbox(sl, title, rx + Cm(0.3), ry + Cm(0.25),
                    rw - Cm(0.5), Cm(0.7),
                    font_size=12, bold=True, color=CYAN)
        add_textbox(sl, desc, rx + Cm(0.3), ry + Cm(1.05),
                    rw - Cm(0.5), Cm(2.7),
                    font_size=10, color=LIGHT_GRAY)

    # 세부 기능 목록
    add_textbox(sl, "세부 기능 목록", Cm(2.0), Cm(8.2), Cm(15), Cm(0.8),
                font_size=14, bold=True, color=WHITE)
    add_rect(sl, Cm(2.0), Cm(8.9), SLIDE_W - Cm(4), Cm(0.04), MID_BLUE)

    funcs = [
        ("상품 광고 카피 생성", "타깃/목적/컨셉 맞춤 상품 광고 카피 자동 생성"),
        ("Alt-text 생성",       "이미지 설명문 및 SEO 키워드 자동 생성"),
        ("금융수준 맞춤 스타일링", "사용자 금융 이해 수준에 맞춘 본문 스타일 변형"),
        ("고객 인어가이드 생성", "고객 대상 맞춤형 안내 문서 자동 생성"),
        ("상품 FAQ 자동 생성",  "상품별 자주 묻는 질문 및 답변 자동 생성"),
        ("이벤트 메시지 생성",  "캠페인·이벤트 홍보 메시지 자동 작성"),
        ("텍스트 스타일 변환",  "격식체/비격식체 등 문체 변환"),
    ]
    for i, (fn, desc) in enumerate(funcs):
        col = i % 2
        row = i // 2
        fx = Cm(2.2) + col * Cm(15.8)
        fy = Cm(9.3) + row * Cm(1.6)
        add_rect(sl, fx, fy, Cm(1.1), Cm(0.65),
                 CYAN if col == 0 else ORANGE)
        add_textbox(sl, f"F{i+1:02d}", fx, fy, Cm(1.1), Cm(0.65),
                    font_size=8, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_textbox(sl, fn, fx + Cm(1.3), fy, Cm(5.5), Cm(0.65),
                    font_size=10, bold=True, color=WHITE)
        add_textbox(sl, desc, fx + Cm(1.3), fy + Cm(0.62), Cm(13), Cm(0.65),
                    font_size=8.5, color=DARK_TEXT)

    add_footer(sl)
    add_slide_number(sl, 5)


# ─── 슬라이드 6 : 문서·콘텐츠 생성 ──────────────────────────────────────────────

def slide_06_document(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)

    add_rect(sl, Cm(1.4), Cm(1.2), Cm(0.3), Cm(1.5), GREEN)
    add_textbox(sl, "문서·콘텐츠 생성", Cm(2.0), Cm(1.2), Cm(18), Cm(1.5),
                font_size=28, bold=True, color=WHITE)
    add_section_tag(sl, "DOCUMENT & CONTENT GENERATION", Cm(2.0), Cm(2.65),
                    bg=GREEN, fg=NAVY)
    add_percent_badge(sl, "77.3%", Cm(25.5), Cm(1.3), GREEN)
    add_textbox(sl, "고성공 영역", Cm(25.0), Cm(2.55), Cm(4.5), Cm(0.6),
                font_size=9, color=GREEN, align=PP_ALIGN.CENTER)

    # 성공 요인
    add_textbox(sl, "성공 요인 분석", Cm(2.0), Cm(3.6), Cm(15), Cm(0.8),
                font_size=14, bold=True, color=GREEN)

    factors = [
        "반복적인 문서 작성 업무를 AI로 전환하여 생산성 대폭 향상",
        "금융 용어·규정 학습된 모델 활용으로 정확도 확보",
        "사람 검토 단계 유지로 오류 위험 관리 가능",
        "다양한 문서 유형(보고서, 공문, 안내문 등)에 범용 적용",
    ]
    for i, f in enumerate(factors):
        fy = Cm(4.4) + i * Cm(0.85)
        add_rect(sl, Cm(2.0), fy + Cm(0.15), Cm(0.55), Cm(0.45), GREEN)
        add_textbox(sl, str(i+1), Cm(2.0), fy + Cm(0.1), Cm(0.55), Cm(0.55),
                    font_size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_textbox(sl, f, Cm(2.8), fy, Cm(28), Cm(0.75),
                    font_size=11, color=LIGHT_GRAY)

    # 세부 기능
    add_rect(sl, Cm(2.0), Cm(8.3), SLIDE_W - Cm(4), Cm(0.04), MID_BLUE)
    add_textbox(sl, "주요 적용 기능", Cm(2.0), Cm(8.5), Cm(15), Cm(0.7),
                font_size=14, bold=True, color=WHITE)

    doc_funcs = [
        ("보고서 자동 작성",   "데이터 기반 경영 보고서·분석 보고서 초안 자동 생성"),
        ("공문·안내문 생성",   "수신자·목적에 맞는 공식 문서 자동 작성"),
        ("계약서 요약·검토",   "금융 계약서 핵심 조항 추출 및 리스크 포인트 요약"),
        ("규정·지침 문서화",   "내부 규정 변경 사항 자동 반영 및 문서 갱신"),
        ("회의록 자동 생성",   "회의 녹취 기반 안건·결정 사항·후속 조치 자동 정리"),
        ("뉴스레터·공지 생성", "내·외부 커뮤니케이션용 뉴스레터 및 공지 자동 작성"),
    ]

    for i, (fn, desc) in enumerate(doc_funcs):
        col = i % 3
        row = i // 3
        fx = Cm(2.0) + col * Cm(10.6)
        fy = Cm(9.4) + row * Cm(3.5)
        fw, fh = Cm(10.0), Cm(3.2)
        add_rect(sl, fx, fy, fw, fh, NAVY2, GREEN, Pt(0.75))
        add_rect(sl, fx, fy, fw, Cm(0.12), GREEN)
        add_textbox(sl, fn, fx + Cm(0.3), fy + Cm(0.25),
                    fw - Cm(0.5), Cm(0.7),
                    font_size=11, bold=True, color=GREEN)
        add_textbox(sl, desc, fx + Cm(0.3), fy + Cm(1.05),
                    fw - Cm(0.5), Cm(1.9),
                    font_size=9.5, color=LIGHT_GRAY)

    add_footer(sl)
    add_slide_number(sl, 6)


# ─── 슬라이드 7 : 상담 지원 ──────────────────────────────────────────────────────

def slide_07_consulting(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)

    ACCENT = RGBColor(0x7B, 0x68, 0xEE)  # medium slate blue

    add_rect(sl, Cm(1.4), Cm(1.2), Cm(0.3), Cm(1.5), ACCENT)
    add_textbox(sl, "상담 지원", Cm(2.0), Cm(1.2), Cm(18), Cm(1.5),
                font_size=28, bold=True, color=WHITE)
    add_section_tag(sl, "CONSULTATION SUPPORT", Cm(2.0), Cm(2.65),
                    bg=ACCENT, fg=WHITE)
    add_percent_badge(sl, "70.0%", Cm(25.5), Cm(1.3), ACCENT)
    add_textbox(sl, "고성공 영역", Cm(25.0), Cm(2.55), Cm(4.5), Cm(0.6),
                font_size=9, color=ACCENT, align=PP_ALIGN.CENTER)

    # 4개 기능 카드
    consulting_cards = [
        ("상담 요약",
         "STT 녹취 요약",
         "콜 녹취 내용을 키워드 추출·불만 요약·질의 유형 자동 분류하여 상담사 업무 시간 70% 절감"),
        ("상담 답변 생성",
         "챗봇 지원 시스템",
         "상품 약관·사업방법서·업무 매뉴얼 기반 상담사 지원 챗봇으로 답변 정확도 향상"),
        ("상담로그 유형화",
         "패턴 자동 분류",
         "전체 상담 로그 내 잠재 유형 자동 도출 및 유형별 건수 계수로 운영 인사이트 확보"),
        ("고객 접점 리포트",
         "터치포인트 분석",
         "콜센터 상담 요약 기반 고객 관심사·이벤트 자동 정의로 마케팅 연계 데이터 생성"),
    ]

    for i, (title, sub, desc) in enumerate(consulting_cards):
        col = i % 2
        row = i // 2
        cx = Cm(2.0) + col * Cm(15.8)
        cy = Cm(3.6) + row * Cm(6.8)
        cw, ch = Cm(15.0), Cm(6.3)
        add_rect(sl, cx, cy, cw, ch, NAVY2, ACCENT, Pt(1.0))
        add_rect(sl, cx, cy, Cm(0.3), ch, ACCENT)
        add_textbox(sl, title, cx + Cm(0.6), cy + Cm(0.25),
                    cw - Cm(0.8), Cm(0.85),
                    font_size=14, bold=True, color=ACCENT)
        add_textbox(sl, sub, cx + Cm(0.6), cy + Cm(1.05),
                    cw - Cm(0.8), Cm(0.6),
                    font_size=10, color=YELLOW, bold=True)
        add_rect(sl, cx + Cm(0.6), cy + Cm(1.65), cw - Cm(0.9), Cm(0.03), MID_BLUE)
        add_textbox(sl, desc, cx + Cm(0.6), cy + Cm(1.8),
                    cw - Cm(0.8), Cm(4.0),
                    font_size=10, color=LIGHT_GRAY)

    add_footer(sl)
    add_slide_number(sl, 7)


# ─── 슬라이드 8 : 투자/리서치 ────────────────────────────────────────────────────

def slide_08_investment(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)

    add_rect(sl, Cm(1.4), Cm(1.2), Cm(0.3), Cm(1.5), YELLOW)
    add_textbox(sl, "투자 / 리서치", Cm(2.0), Cm(1.2), Cm(18), Cm(1.5),
                font_size=28, bold=True, color=WHITE)
    add_section_tag(sl, "INVESTMENT & RESEARCH", Cm(2.0), Cm(2.65),
                    bg=YELLOW, fg=NAVY)
    add_percent_badge(sl, "42.9%", Cm(25.5), Cm(1.3), YELLOW)
    add_textbox(sl, "중간 성공 영역", Cm(24.8), Cm(2.55), Cm(4.8), Cm(0.6),
                font_size=9, color=YELLOW, align=PP_ALIGN.CENTER)

    # 도전 과제 박스
    add_rect(sl, Cm(2.0), Cm(3.5), SLIDE_W - Cm(4), Cm(2.5),
             RGBColor(0x1A, 0x2A, 0x10), RGBColor(0xFF, 0xD7, 0x00), Pt(1.0))
    add_textbox(sl, "⚠  낮은 성공률 원인",
                Cm(2.4), Cm(3.7), Cm(20), Cm(0.75),
                font_size=12, bold=True, color=YELLOW)
    challenges = "규제 민감도 높음 (투자 권유·소비자 보호 관련) · 할루시네이션 위험 직접적 금전 손실 연결 · 고난이도 전문성 검증 체계 미비"
    add_textbox(sl, challenges, Cm(2.4), Cm(4.5), SLIDE_W - Cm(5), Cm(1.3),
                font_size=10, color=LIGHT_GRAY)

    # 세부 기능
    invest_funcs = [
        ("지식 공유", YELLOW,
         "직원용 보고서 및 정보 제공\n상품-하우스뷰/시황 연계 질의응답\n내부 투자 교육 자료 자동화"),
        ("상품 제안", ORANGE,
         "포트폴리오 고객용 눈높이 설명서\n스크립트 자동 작성\n맞춤형 투자 제안서 생성"),
        ("뉴스 검색/요약", CYAN,
         "실시간 금융 뉴스 수집 및 요약\n주요 이슈 자동 브리핑\n시장 영향 분석 지원"),
        ("기업 분석 보고서", GREEN,
         "재무제표 기반 자동 분석\n기업 리포트 초안 생성\n비교 분석 자동화"),
        ("금융 키워드 선별", YELLOW,
         "시장 트렌드 키워드 자동 추출\n투자 테마 분류\n리포트 색인 자동화"),
    ]

    for i, (fn, color, desc) in enumerate(invest_funcs):
        col = i % 3
        row = i // 3
        fx = Cm(2.0) + col * Cm(10.6)
        fy = Cm(6.5) + row * Cm(5.5)
        fw, fh = Cm(10.0), Cm(5.1)
        add_rect(sl, fx, fy, fw, fh, NAVY2, color, Pt(0.75))
        add_rect(sl, fx, fy, fw, Cm(0.12), color)
        add_textbox(sl, fn, fx + Cm(0.3), fy + Cm(0.25),
                    fw - Cm(0.5), Cm(0.7),
                    font_size=12, bold=True, color=color)
        add_textbox(sl, desc, fx + Cm(0.3), fy + Cm(1.1),
                    fw - Cm(0.5), Cm(3.7),
                    font_size=10, color=LIGHT_GRAY)

    add_footer(sl)
    add_slide_number(sl, 8)


# ─── 슬라이드 9 : 리스크/보안 ────────────────────────────────────────────────────

def slide_09_risk(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)

    add_rect(sl, Cm(1.4), Cm(1.2), Cm(0.3), Cm(1.5), ORANGE)
    add_textbox(sl, "리스크 / 보안", Cm(2.0), Cm(1.2), Cm(18), Cm(1.5),
                font_size=28, bold=True, color=WHITE)
    add_section_tag(sl, "RISK & SECURITY", Cm(2.0), Cm(2.65),
                    bg=ORANGE, fg=WHITE)
    add_percent_badge(sl, "37.5%", Cm(25.5), Cm(1.3), ORANGE)
    add_textbox(sl, "개선 필요 영역", Cm(24.8), Cm(2.55), Cm(4.8), Cm(0.6),
                font_size=9, color=ORANGE, align=PP_ALIGN.CENTER)

    # 낮은 성공률 원인
    add_rect(sl, Cm(2.0), Cm(3.5), SLIDE_W - Cm(4), Cm(2.8),
             RGBColor(0x2A, 0x10, 0x10), ORANGE, Pt(1.0))
    add_textbox(sl, "낮은 성공률 배경",
                Cm(2.4), Cm(3.7), Cm(20), Cm(0.75),
                font_size=12, bold=True, color=ORANGE)
    bg_text = ("규제 당국의 자동화 의사결정 제한 (금융감독원 가이드라인)  ·  "
               "오탐(False Positive) 발생 시 심각한 컴플라이언스 리스크  ·  "
               "보안 데이터의 AI 모델 외부 전송 우려  ·  "
               "기존 레거시 리스크 시스템과의 통합 난이도 높음")
    add_textbox(sl, bg_text, Cm(2.4), Cm(4.55), SLIDE_W - Cm(5), Cm(1.5),
                font_size=10, color=LIGHT_GRAY)

    # 세부 기능 3개
    risk_funcs = [
        ("사내 핵심 내용 추출", ORANGE,
         "내부 문서에서 리스크 관련 핵심 내용 자동 추출\n이상 징후 탐지를 위한 텍스트 분석\n정책 위반 패턴 자동 감지",
         "적용 단계: 내부 문서 검토 보조"),
        ("성문 일치 검증", RGBColor(0xFF, 0x99, 0x55),
         "계약서·약관의 규정 준수 여부 자동 검토\n법적 요구사항과의 불일치 자동 탐지\n개정 전후 변경 내용 자동 비교",
         "적용 단계: 컴플라이언스 검토 보조"),
        ("고객 맞춤형 서류 안내", YELLOW,
         "업무/사유별 제출 서류 자동 안내\n개인별 필요 서류 목록 맞춤 제공\n서류 누락 방지 및 재방문율 감소",
         "적용 단계: 고객 응대 지원"),
    ]

    for i, (fn, color, desc, stage) in enumerate(risk_funcs):
        fx = Cm(2.0) + i * Cm(10.6)
        fy = Cm(6.8)
        fw, fh = Cm(10.0), Cm(9.5)
        add_rect(sl, fx, fy, fw, fh, NAVY2, color, Pt(0.75))
        add_rect(sl, fx, fy, fw, Cm(0.12), color)
        add_textbox(sl, fn, fx + Cm(0.3), fy + Cm(0.25),
                    fw - Cm(0.5), Cm(0.8),
                    font_size=12, bold=True, color=color)
        add_textbox(sl, desc, fx + Cm(0.3), fy + Cm(1.15),
                    fw - Cm(0.5), Cm(6.5),
                    font_size=10, color=LIGHT_GRAY)
        add_rect(sl, fx + Cm(0.3), fy + fh - Cm(1.3),
                 fw - Cm(0.5), Cm(0.04), MID_BLUE)
        add_textbox(sl, stage, fx + Cm(0.3), fy + fh - Cm(1.2),
                    fw - Cm(0.5), Cm(0.6),
                    font_size=9, color=color, bold=True)

    add_footer(sl)
    add_slide_number(sl, 9)


# ─── 슬라이드 10 : 내부 운영 지원 ────────────────────────────────────────────────

def slide_10_internal(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)

    ACCENT2 = RGBColor(0xA0, 0x82, 0xFF)

    add_rect(sl, Cm(1.4), Cm(1.2), Cm(0.3), Cm(1.5), ACCENT2)
    add_textbox(sl, "내부 운영 지원", Cm(2.0), Cm(1.2), Cm(18), Cm(1.5),
                font_size=28, bold=True, color=WHITE)
    add_section_tag(sl, "INTERNAL OPERATIONS SUPPORT", Cm(2.0), Cm(2.65),
                    bg=ACCENT2, fg=WHITE)
    add_percent_badge(sl, "33.3%", Cm(25.5), Cm(1.3), ACCENT2)
    add_textbox(sl, "개선 필요 영역", Cm(24.8), Cm(2.55), Cm(4.8), Cm(0.6),
                font_size=9, color=ACCENT2, align=PP_ALIGN.CENTER)

    # 원인 분석
    add_rect(sl, Cm(2.0), Cm(3.5), SLIDE_W - Cm(4), Cm(2.5),
             RGBColor(0x15, 0x10, 0x2A), ACCENT2, Pt(1.0))
    add_textbox(sl, "성공률 저하 주요 원인",
                Cm(2.4), Cm(3.7), Cm(20), Cm(0.75),
                font_size=12, bold=True, color=ACCENT2)
    cause_text = ("완전 자동화 추진 과정에서 예외 처리 로직 복잡도 과소평가  ·  "
                  "레거시 시스템 연동 난이도  ·  "
                  "변화 관리 및 직원 수용도 부족  ·  "
                  "명확한 ROI 측정 지표 부재")
    add_textbox(sl, cause_text, Cm(2.4), Cm(4.55), SLIDE_W - Cm(5), Cm(1.2),
                font_size=10, color=LIGHT_GRAY)

    # 3개 기능 카드
    internal_funcs = [
        ("내부 문서 자동 분류", ACCENT2,
         "문서 유형·부서·중요도 기준 자동 분류\n태그 자동 부여 및 검색 최적화\n문서 생애주기 관리 자동화",
         "예상 효과: 문서 검색 시간 60% 절감"),
        ("업무 프로세스 자동화", RGBColor(0xC0, 0xA0, 0xFF),
         "반복 업무 워크플로우 자동화\n승인 프로세스 지능형 라우팅\n예외 상황 자동 감지 및 알림",
         "예상 효과: 반복 업무 처리 시간 40% 절감"),
        ("내부 지식 Q&A", CYAN,
         "사내 지식 베이스 기반 질의응답\n직원 온보딩 지원 자동화\n업무 매뉴얼 접근성 향상",
         "예상 효과: 내부 문의 건수 35% 감소"),
    ]

    for i, (fn, color, desc, effect) in enumerate(internal_funcs):
        fx = Cm(2.0) + i * Cm(10.6)
        fy = Cm(6.5)
        fw, fh = Cm(10.0), Cm(9.8)
        add_rect(sl, fx, fy, fw, fh, NAVY2, color, Pt(0.75))
        add_rect(sl, fx, fy, fw, Cm(0.12), color)
        add_textbox(sl, fn, fx + Cm(0.3), fy + Cm(0.25),
                    fw - Cm(0.5), Cm(0.8),
                    font_size=12, bold=True, color=color)
        add_textbox(sl, desc, fx + Cm(0.3), fy + Cm(1.15),
                    fw - Cm(0.5), Cm(6.8),
                    font_size=10, color=LIGHT_GRAY)
        add_rect(sl, fx + Cm(0.3), fy + fh - Cm(1.4),
                 fw - Cm(0.5), Cm(0.04), color)
        add_textbox(sl, effect, fx + Cm(0.3), fy + fh - Cm(1.3),
                    fw - Cm(0.5), Cm(0.6),
                    font_size=9, color=color, bold=True)

    add_footer(sl)
    add_slide_number(sl, 10)


# ─── 슬라이드 11 : 시사점 및 전략 ────────────────────────────────────────────────

def slide_11_insights(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)

    add_rect(sl, Cm(1.4), Cm(1.2), Cm(0.3), Cm(1.5), CYAN)
    add_textbox(sl, "시사점 및 전략적 제언", Cm(2.0), Cm(1.2), Cm(20), Cm(1.5),
                font_size=28, bold=True, color=WHITE)
    add_section_tag(sl, "INSIGHTS & STRATEGIC RECOMMENDATIONS", Cm(2.0), Cm(2.65))

    # 3개 섹션 카드
    sections = [
        ("성공의 공통점", GREEN, [
            "명확한 KPI 설정 (클릭률·처리 시간 등 정량 지표)",
            "인간 검토 단계 반드시 포함 (Human-in-the-Loop)",
            "좁고 명확한 태스크 정의 (범용 AI 지양)",
            "파일럿 → 부서 전개 → 전사 확산 단계적 접근",
            "현업 담당자 적극 참여 및 피드백 루프 운영",
        ]),
        ("실패 패턴", ORANGE, [
            "완전 자동화 시도 (예외 처리 간과)",
            "규제 민감 영역 우선 시도 (컴플라이언스 리스크)",
            "KPI 미설정 상태로 도입 후 효과 측정 불가",
            "IT 주도 추진 (현업 수요 미반영)",
            "레거시 시스템 통합 난이도 과소평가",
        ]),
        ("도입 권장 순서", CYAN, [
            "① 고객/마케팅 (규제 낮음·효과 즉시)",
            "② 상담 지원 (효율성 높음·품질 관리 용이)",
            "③ 문서·콘텐츠 생성 (범용성 높음)",
            "④ 투자/리서치 (전문성 검증 체계 선구축)",
            "⑤ 리스크/보안 (규제 정비 후 단계 도입)",
        ]),
    ]

    for i, (title, color, bullets) in enumerate(sections):
        sx = Cm(2.0) + i * Cm(10.6)
        sy = Cm(3.6)
        sw, sh = Cm(10.0), Cm(12.0)
        add_rect(sl, sx, sy, sw, sh, NAVY2, color, Pt(1.0))
        add_rect(sl, sx, sy, sw, Cm(0.12), color)
        add_textbox(sl, title, sx + Cm(0.3), sy + Cm(0.25),
                    sw - Cm(0.5), Cm(0.8),
                    font_size=13, bold=True, color=color)
        add_rect(sl, sx + Cm(0.3), sy + Cm(1.1),
                 sw - Cm(0.5), Cm(0.03), MID_BLUE)
        for j, b in enumerate(bullets):
            add_textbox(sl, f"•  {b}",
                        sx + Cm(0.3), sy + Cm(1.3) + j * Cm(2.0),
                        sw - Cm(0.5), Cm(1.8),
                        font_size=10, color=LIGHT_GRAY)

    add_footer(sl)
    add_slide_number(sl, 11)


# ─── 슬라이드 12 : 하나은행 적용 방향 ───────────────────────────────────────────

def slide_12_hana(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)

    add_rect(sl, Cm(1.4), Cm(1.2), Cm(0.3), Cm(1.5), GREEN)
    add_textbox(sl, "하나은행 적용 방향", Cm(2.0), Cm(1.2), Cm(20), Cm(1.5),
                font_size=28, bold=True, color=WHITE)
    add_section_tag(sl, "HANA BANK IMPLEMENTATION ROADMAP", Cm(2.0), Cm(2.65),
                    bg=GREEN, fg=NAVY)

    # 비정형 → AI Agent 플로우
    add_textbox(sl, "비정형 데이터 자산화 플랫폼 → AI Agent 연계 구조",
                Cm(2.0), Cm(3.4), Cm(28), Cm(0.8),
                font_size=13, bold=True, color=CYAN)

    # 플로우 박스 5개
    flow_items = [
        ("STEP 1", "비정형 데이터\n수집·정제", CYAN,
         "상담 녹취, 문서,\n뉴스, 보고서"),
        ("STEP 2", "AI 임베딩\n및 지식화", GREEN,
         "벡터 DB 구축\nRAG 파이프라인"),
        ("STEP 3", "AI Agent\n플랫폼 구축", YELLOW,
         "오케스트레이션\n레이어 설계"),
        ("STEP 4", "파일럿\n영역 선정", ORANGE,
         "고객/마케팅\n상담 지원 우선"),
        ("STEP 5", "전사 확산\n및 고도화", RGBColor(0xA0, 0x82, 0xFF),
         "성과 측정 후\n단계적 확대"),
    ]

    for i, (step, title, color, desc) in enumerate(flow_items):
        fx = Cm(2.0) + i * Cm(6.2)
        fy = Cm(4.4)
        fw, fh = Cm(5.6), Cm(5.5)
        add_rect(sl, fx, fy, fw, fh, NAVY2, color, Pt(1.0))
        add_rect(sl, fx, fy, fw, Cm(0.12), color)
        # 스텝 배지
        add_rect(sl, fx + Cm(0.3), fy + Cm(0.3), Cm(2.0), Cm(0.55), color)
        add_textbox(sl, step, fx + Cm(0.3), fy + Cm(0.3), Cm(2.0), Cm(0.55),
                    font_size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_textbox(sl, title, fx + Cm(0.3), fy + Cm(1.1),
                    fw - Cm(0.5), Cm(1.4),
                    font_size=12, bold=True, color=color)
        add_textbox(sl, desc, fx + Cm(0.3), fy + Cm(2.7),
                    fw - Cm(0.5), Cm(2.5),
                    font_size=10, color=LIGHT_GRAY)
        # 화살표 (마지막 제외)
        if i < 4:
            add_textbox(sl, "▶", fx + fw + Cm(0.05), fy + Cm(2.3),
                        Cm(0.6), Cm(0.8),
                        font_size=14, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    # 기대 효과 표
    add_textbox(sl, "단계별 기대 효과 및 목표 KPI",
                Cm(2.0), Cm(10.2), Cm(28), Cm(0.7),
                font_size=13, bold=True, color=WHITE)
    add_rect(sl, Cm(2.0), Cm(10.9), SLIDE_W - Cm(4), Cm(0.04), MID_BLUE)

    kpi_items = [
        ("고객/마케팅 Agent", "콘텐츠 생산량 10배↑", "3개월 내 본사업 전환", CYAN),
        ("상담 지원 Agent",   "상담 처리 시간 70%↓",  "콜센터 파일럿 6개월", RGBColor(0x7B, 0x68, 0xEE)),
        ("문서 생성 Agent",   "보고서 작성 시간 80%↓", "전사 확산 12개월",   GREEN),
        ("리스크 Agent",      "컴플라이언스 오류 50%↓","규제 검토 후 도입",  ORANGE),
    ]

    for i, (area, effect, timeline, color) in enumerate(kpi_items):
        col = i % 2
        row = i // 2
        kx = Cm(2.0) + col * Cm(16.0)
        ky = Cm(11.4) + row * Cm(2.8)
        kw, kh = Cm(15.0), Cm(2.4)
        add_rect(sl, kx, ky, kw, kh, NAVY2, color, Pt(0.5))
        add_rect(sl, kx, ky, Cm(0.25), kh, color)
        add_textbox(sl, area, kx + Cm(0.5), ky + Cm(0.15),
                    Cm(5.5), Cm(0.7),
                    font_size=11, bold=True, color=color)
        add_textbox(sl, effect, kx + Cm(0.5), ky + Cm(0.85),
                    Cm(8), Cm(0.65),
                    font_size=10, color=WHITE)
        add_textbox(sl, timeline, kx + Cm(9.5), ky + Cm(0.5),
                    Cm(5.0), Cm(0.65),
                    font_size=9, color=DARK_TEXT, align=PP_ALIGN.RIGHT)

    add_footer(sl)
    add_slide_number(sl, 12)


# ─── 메인 실행 ────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    print("슬라이드 생성 중...")
    slide_01_cover(prs)
    print("  [1/12] 표지")
    slide_02_toc(prs)
    print("  [2/12] 목차")
    slide_03_overview(prs)
    print("  [3/12] 조사 개요")
    slide_04_chart(prs)
    print("  [4/12] 성공률 차트")
    slide_05_marketing(prs)
    print("  [5/12] 고객/마케팅")
    slide_06_document(prs)
    print("  [6/12] 문서·콘텐츠 생성")
    slide_07_consulting(prs)
    print("  [7/12] 상담 지원")
    slide_08_investment(prs)
    print("  [8/12] 투자/리서치")
    slide_09_risk(prs)
    print("  [9/12] 리스크/보안")
    slide_10_internal(prs)
    print("  [10/12] 내부 운영 지원")
    slide_11_insights(prs)
    print("  [11/12] 시사점")
    slide_12_hana(prs)
    print("  [12/12] 하나은행 적용 방향")

    output_path = "/tmp/금융권_AI_Agent_시장조사_보고서.pptx"
    prs.save(output_path)
    print(f"\n✓ 파일 저장 완료: {output_path}")

    import os
    size = os.path.getsize(output_path)
    print(f"  파일 크기: {size / 1024:.1f} KB")
    print(f"  슬라이드 수: {len(prs.slides)}")


if __name__ == "__main__":
    main()
