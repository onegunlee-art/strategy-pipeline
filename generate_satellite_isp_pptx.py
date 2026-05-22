"""
위성영상 ISP 수주전략 보고서 — 4장 슬라이드 PPT 생성
입력 이미지에 표시된 글자를 그대로 재구성.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from copy import deepcopy

# ---- 컬러 ----
COLOR_HEADER_BG = RGBColor(0xE7, 0xEB, 0xF0)
COLOR_BORDER = RGBColor(0x9E, 0xA7, 0xB0)
COLOR_TEXT = RGBColor(0x1F, 0x29, 0x33)
COLOR_TITLE = RGBColor(0x21, 0x3A, 0x6B)
COLOR_RED_TAG = RGBColor(0xC0, 0x1B, 0x2E)
COLOR_BLUE_TAG = RGBColor(0x1F, 0x4E, 0x9C)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

def add_textbox(slide, left, top, width, height, text, *,
                font_size=11, bold=False, color=COLOR_TEXT,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(50000)
    tf.margin_right = Emu(50000)
    tf.margin_top = Emu(30000)
    tf.margin_bottom = Emu(30000)
    tf.vertical_anchor = anchor

    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = '맑은 고딕'
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb

def add_cell(slide, left, top, width, height, text, *,
             header=False, font_size=10.5, align=PP_ALIGN.LEFT,
             bold=None, color=None, fill=None,
             anchor=MSO_ANCHOR.MIDDLE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.color.rgb = COLOR_BORDER
    shape.line.width = Pt(0.5)

    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    elif header:
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_HEADER_BG
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_WHITE

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(70000)
    tf.margin_right = Emu(70000)
    tf.margin_top = Emu(40000)
    tf.margin_bottom = Emu(40000)
    tf.vertical_anchor = anchor

    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.18
        run = p.add_run()
        run.text = line
        run.font.name = '맑은 고딕'
        run.font.size = Pt(font_size)
        run.font.bold = bold if bold is not None else header
        run.font.color.rgb = color if color is not None else COLOR_TEXT
    return shape

def add_tag(slide, left, top, width, height, text, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    tf = shape.text_frame
    tf.margin_left = Emu(30000); tf.margin_right = Emu(30000)
    tf.margin_top = Emu(20000); tf.margin_bottom = Emu(20000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = '맑은 고딕'
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = COLOR_WHITE
    return shape

def add_title(slide, text):
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                text, font_size=24, bold=True, color=COLOR_TITLE)

# ============================================================
# Slide 1 — 사업내용 + 추진 방향 (page 2/6)
# ============================================================
s = prs.slides.add_slide(BLANK)

# Left table
TL = Inches(0.5)
TT = Inches(1.2)
COL1_W = Inches(1.4)
COL2_W = Inches(6.5)
ROW1_H = Inches(0.45)

# Header row
add_cell(s, TL, TT, COL1_W, ROW1_H, '구  분', header=True, bold=True, align=PP_ALIGN.CENTER)
add_cell(s, TL + COL1_W, TT, COL2_W, ROW1_H, '사업내용', header=True, bold=True, align=PP_ALIGN.CENTER)

# Row 1 — 소프트웨어 개발
r1_h = Inches(1.5)
y = TT + ROW1_H
add_cell(s, TL, y, COL1_W, r1_h, '소프트웨어\n개발', align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, bold=True)
sw_text = (
    "○ 통합 수신·처리시스템 개발\n"
    "○ 영상관리 시스템 개발 / 대외배포시스템 개발\n"
    "○ 통합 촬영계획시스템 개발\n"
    "○ UI 설계, 데이터 이관 및 관련 도구 개발"
)
add_cell(s, TL + COL1_W, y, COL2_W, r1_h, sw_text)

# Row 2 — 인프라 도입
y += r1_h
r2_h = Inches(1.15)
add_cell(s, TL, y, COL1_W, r2_h, '인프라 도입', align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, bold=True)
inf_text = (
    "○ 영상저장 인프라 구축 / 기반 아키텍처 구축\n"
    "○ 네트워크 아키텍처 구축 / 보안 아키텍처 구축\n"
    "○ VDI 시스템 구축 / IT 운영 관제시스템 구축"
)
add_cell(s, TL + COL1_W, y, COL2_W, r2_h, inf_text)

# Row 3 — 기반환경 구축
y += r2_h
r3_h = Inches(2.5)
add_cell(s, TL, y, COL1_W, r3_h, '기반환경 구축', align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, bold=True)
env_text = (
    "○ 통합 운영실 인테리어 공사\n"
    "○ 전산실 컨테인먼트 구성 및 내부 마감\n"
    "○ 공조 체계 구축, 공조 배관 및 누수 감시 체계 구축\n"
    "○ UPS·배터리 구성, 전산실 전원 공사 및 케이블 포설\n"
    "○ 전력용 Tray·부스덕트 및 모니터링 체계 구축\n"
    "○ 네트워크 케이블·Tray 공사\n"
    "○ CCTV 구축 및 DCIM 체계 구축"
)
add_cell(s, TL + COL1_W, y, COL2_W, r3_h, env_text)

# Right — 추진 방향
RL = Inches(8.8)
add_textbox(s, RL, Inches(0.9), Inches(4.2), Inches(0.7),
            '추진 방향', font_size=28, bold=True, color=COLOR_TITLE, align=PP_ALIGN.CENTER)

bullets = [
    "위성 확대와 대용량 영상 데이터를 고려한 구조적\n유연성·효율성 확보 및 최신기술 적용",
    "신속한 위성영상 처리·판독과 외부 위협으로부터\n안전한 데이터 저장 체계 마련",
    "군·유관기관 등 활용 가능한 위성영상 대외 배포\n시스템 구축",
    "비상 상황 대비 업무 연속성 보장을 위한 재해복구\n및 백업체계 구축",
]
by = Inches(2.0)
for b in bullets:
    add_textbox(s, RL, by, Inches(4.4), Inches(1.0), '• ' + b,
                font_size=12, line_spacing=1.3)
    by += Inches(1.05)

# ============================================================
# Slide 2 — 1. 고객/사업 현황 분석 (page 3/6)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_title(s, '1. 고객/사업 현황 분석')

TL = Inches(0.5)
TT = Inches(1.3)
COL1_W = Inches(2.2)
COL2_W = Inches(10.4)
HDR_H = Inches(0.45)

add_cell(s, TL, TT, COL1_W, HDR_H, '과업범위', header=True, bold=True, align=PP_ALIGN.CENTER)
add_cell(s, TL + COL1_W, TT, COL2_W, HDR_H, '고객니즈/페인포인트', header=True, bold=True, align=PP_ALIGN.CENTER)

y = TT + HDR_H
# Row 1 — 소프트웨어 개발
r_h = Inches(1.35)
add_cell(s, TL, y, COL1_W, r_h, '소프트웨어 개발', align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, bold=True)
sw_pp = (
    "• [데이터 표준화] 현재 각기 다른 표준의 위성영상 데이터 표준화 니즈\n"
    "• [데이터 이관] 서울 센터 peta 단위 데이터를 제주로 안정적/신속한 이관 원함\n"
    "• [시스템 통합] Silo한 위성영상 시스템을 통합하여 구조적 유연성 효율성 확보 원함"
)
add_cell(s, TL + COL1_W, y, COL2_W, r_h, sw_pp)

# Row 2 — 인프라 도입 및 환경 구축
y += r_h
r_h2 = Inches(1.85)
add_cell(s, TL, y, COL1_W, r_h2, '인프라 도입 및\n환경 구축', align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, bold=True)
inf_pp = (
    "• [VDI 환경 구축] 서울-제주 실시간 판독 분석 가능한 VDI 환경 필요\n"
    "• [DC 구축/운영] 국자원, 항우연 등 유사 사례 기반 안정적 구축 및 무중단 운영 방안 요청\n"
    "• [백업체계 구축] 비상상황에도 업무 연속성 보장하는 끊김 없는 업무 체계 필요\n"
    "• [운영/관제] 인력 감소/업무 증대로 AI 등 활용 통한 자동화, 효율화 필요"
)
add_cell(s, TL + COL1_W, y, COL2_W, r_h2, inf_pp)

# Row 3 — 기타
y += r_h2
r_h3 = Inches(1.85)
add_cell(s, TL, y, COL1_W, r_h3, '기타', align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, bold=True)
etc_pp = (
    "• [업무 자동화] 보고 업무 지원 등 “일부” 업무 영역 자동화 관심(판독업무 X)\n"
    "• [미래 모델 제시] 미국 NCG 등 향후 선진화된 사례 도입에 관심(대외 배포 등)\n"
    "• [사업관리] 4개년 대형사업에 맞는 책임있는 사업관리, 이슈/리스크 관리 필요\n"
    "• [gray영역 존재] ISP RFP 상 명확화 되지 않은 영역 다수(향후 구축, 기술협상 시 논의 원함)"
)
add_cell(s, TL + COL1_W, y, COL2_W, r_h3, etc_pp)

y += r_h3
add_textbox(s, TL, y + Inches(0.15), Inches(12), Inches(0.4),
            '* 평가 위원 현황: 고객사 내부 유관 부서 실무 담당자 8인으로 구성 예상, 인프라/운영/기타 파트 등',
            font_size=11)

# ============================================================
# Slide 3 — 2. 경쟁 현황 분석 (page 4/6)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_title(s, '2. 경쟁 현황 분석')

TL = Inches(0.5)
TT = Inches(1.3)
COL1_W = Inches(2.0)
COL2_W = Inches(5.3)
COL3_W = Inches(5.3)
HDR_H = Inches(0.45)

add_cell(s, TL, TT, COL1_W, HDR_H, '과업범위', header=True, bold=True, align=PP_ALIGN.CENTER)
add_cell(s, TL + COL1_W, TT, COL2_W, HDR_H, 'KT', header=True, bold=True, align=PP_ALIGN.CENTER, font_size=12)
add_cell(s, TL + COL1_W + COL2_W, TT, COL3_W, HDR_H, 'LG CNS', header=True, bold=True, align=PP_ALIGN.CENTER, font_size=12)

y = TT + HDR_H
# Row 1 — 소프트웨어 개발
r_h = Inches(1.8)
add_cell(s, TL, y, COL1_W, r_h, '소프트웨어\n개발', align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, bold=True)

kt_sw = (
    "컨소시엄 수행 예상 영역\n"
    "협력 업체: 컨텍, 한컴 인터페이스(현 구축/운용사)\n"
    "• 개발 관련 고객사 내부 상대적 인지도/이해도 부족 인식\n"
    "• KT SAT 역량 활용 가능"
)
add_cell(s, TL + COL1_W, y, COL2_W, r_h, kt_sw, font_size=10)
add_tag(s, TL + COL1_W + COL2_W - Inches(0.9), y + Inches(0.08), Inches(0.78), Inches(0.3),
        '열위', COLOR_RED_TAG)

lg_sw = (
    "컨소시엄 수행 예상 영역\n"
    "협력 예상 업체: 컨텍+세트렉아이, 한컴 인터페이스\n"
    "• ISP 사업 수행/ G시스템 전신 구축 기반 시스템\n"
    "  이해도 높음"
)
add_cell(s, TL + COL1_W + COL2_W, y, COL3_W, r_h, lg_sw, font_size=10)
add_tag(s, TL + COL1_W + COL2_W + COL3_W - Inches(0.9), y + Inches(0.08), Inches(0.78), Inches(0.3),
        '우위', COLOR_BLUE_TAG)

# Row 2 — 인프라 도입 및 환경 구축
y += r_h
r_h2 = Inches(2.0)
add_cell(s, TL, y, COL1_W, r_h2, '인프라\n도입 및\n환경 구축', align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, bold=True)

kt_inf = (
    "주관사 수행 예상 영역\n"
    "• DC 사업 분야에 대한 긍정적 인식\n"
    "    -  국자원 사업 등 공공 대형 레퍼런스\n"
    "• 수요처 백업망 현 구축/운영사\n"
    "    -  백업망 증설, 서울-제주간 추가 네트워크 이점"
)
add_cell(s, TL + COL1_W, y, COL2_W, r_h2, kt_inf, font_size=10)
add_tag(s, TL + COL1_W + COL2_W - Inches(0.9), y + Inches(0.08), Inches(0.78), Inches(0.3),
        '우위', COLOR_BLUE_TAG)

lg_inf = (
    "주관사 수행 예상 영역\n"
    "• 최근 삼성전자 등 대규모 사업 수주, 내부 구축 건\n"
    "  등을 통한 DC 관련 규모/ 레퍼런스 강조 예상\n"
    "• 운영/관제 관련 엑사원 등 자사 AI 모델 활용한 추\n"
    "  가 제안 예상"
)
add_cell(s, TL + COL1_W + COL2_W, y, COL3_W, r_h2, lg_inf, font_size=10)

# Row 3 — 기타
y += r_h2
r_h3 = Inches(0.85)
add_cell(s, TL, y, COL1_W, r_h3, '기타', align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, bold=True)
add_cell(s, TL + COL1_W, y, COL2_W, r_h3, '', font_size=10)
lg_etc = "• ISP 사업 수행 시 기간연장, 비용 증가 이슈로 인\n  한 고객 불만 제기"
add_cell(s, TL + COL1_W + COL2_W, y, COL3_W, r_h3, lg_etc, font_size=10)

# ============================================================
# Slide 4 — 3. 차별화 아이템 도출 (page 5/6)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_title(s, '3. 차별화 아이템 도출')

TL = Inches(0.5)
TT = Inches(1.3)
COL1_W = Inches(2.2)
COL2_W = Inches(10.4)
HDR_H = Inches(0.45)

add_cell(s, TL, TT, COL1_W, HDR_H, '과업범위', header=True, bold=True, align=PP_ALIGN.CENTER)
add_cell(s, TL + COL1_W, TT, COL2_W, HDR_H, '고객니즈/페인포인트', header=True, bold=True, align=PP_ALIGN.CENTER)

y = TT + HDR_H
# Row 1 — 소프트웨어 개발 (추가 발굴 필요 태그)
r_h = Inches(1.55)
add_cell(s, TL, y, COL1_W, r_h, '소프트웨어 개발', align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, bold=True)
add_tag(s, TL + Inches(0.35), y + Inches(0.75), Inches(1.5), Inches(0.35),
        '추가 발굴 필요', COLOR_RED_TAG)

sw_diff = (
    "• [통합 시스템 개발] 현 시스템 구축 운용경험 기반 통합 시스템 개발로 전처리~촬영계획까지\n"
    "  원활한 연계 보장\n"
    "• [데이터 이관] 실사용 환경 유사 PoC통한 네트워크 방식 무중단/신속한 데이터 이관\n"
    "  -100GB급 서울-제주 연결 전용회선 제공"
)
add_cell(s, TL + COL1_W, y, COL2_W, r_h, sw_diff)

# Row 2 — 인프라 도입 및 환경 구축
y += r_h
r_h2 = Inches(2.1)
add_cell(s, TL, y, COL1_W, r_h2, '인프라 도입 및\n환경 구축', align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, bold=True)
inf_diff = (
    "• [DC 구축/운영] 국자원 등 LL기반 공공 대형센터 구축 노하우 이식 통한 안정적 DC 구축\n"
    "• [백업체계 확보] 서울-제주간 3중화 가능 100기가급 백업망 제공/국사 이중화 등으로 추가\n"
    "    백업 네트워크 안정화\n"
    "• [운영/관제] 국자제 검증 솔루션(엔서블) 및 AI ops 통한 운영관제 자동화\n"
    "• [운영/관제] 최근접 지사 내 출동 기반 KT+협력사 직원 현장 대응 지원 체계"
)
add_cell(s, TL + COL1_W, y, COL2_W, r_h2, inf_diff)

# Row 3 — 기타 (추가 발굴 필요 태그)
y += r_h2
r_h3 = Inches(1.55)
add_cell(s, TL, y, COL1_W, r_h3, '기타', align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, bold=True)
add_tag(s, TL + Inches(0.35), y + Inches(0.55), Inches(1.5), Inches(0.35),
        '추가 발굴 필요', COLOR_RED_TAG)

etc_diff = (
    "• [업무 자동화] AI 적용 업무 영역 식별 및 도입 컨설팅 제공\n"
    "• [미래 모델 제시] 향후 NGA 등 선진 모델 기반 장기 로드맵 제시 및 위성 전문가 컨설팅\n"
    "• [사업 관리] 대형,유사 사업 경험 기반 책임 있는 사업/리스크 관리 방법론 제시"
)
add_cell(s, TL + COL1_W, y, COL2_W, r_h3, etc_diff)

# ============================================================
out = '위성영상_ISP_수주전략.pptx'
prs.save(out)
print(f'saved: {out}')
