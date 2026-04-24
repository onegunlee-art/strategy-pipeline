"""
PPT 빌더 — python-pptx 기반 Layer 1/2/3 export
단일 소스에서 선택 레이어만 export
"""
from __future__ import annotations

import json
import re
from datetime import date
from pathlib import Path

import anthropic
from dotenv import load_dotenv

load_dotenv()
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

from .narrative_arc import get_layer_slides

ROOT = Path(__file__).parent.parent.parent

KT_PRIMARY = "003087"
KT_RED = "E60012"
KT_LIGHT = "F5F7FA"
KT_DARK = "1A1A2E"
KT_GRAY = "888888"
KT_WHITE = "FFFFFF"
KT_ACCENT = "0066CC"


def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_textbox(slide, left, top, width, height, text, font_name="맑은 고딕",
                 font_size=14, bold=False, color_hex=KT_DARK,
                 align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)


def _fill(shape, hex_str: str):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex_str)
    shape.line.fill.background()


def _build_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(5143500)
    return prs


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _header(slide, title: str, W: int, H: int, score_tag: str = "", score_target: str = ""):
    bar = slide.shapes.add_shape(1, 0, 0, Emu(W), Emu(int(H * 0.12)))
    _fill(bar, KT_PRIMARY)

    _add_textbox(slide, int(W*0.04), int(H*0.02), int(W*0.72), int(H*0.1),
                 title, font_size=22, bold=True, color_hex=KT_WHITE)

    if score_tag:
        tag_text = f"📊 {score_tag}"
        if score_target:
            tag_text += f"  {score_target}"
        _add_textbox(slide, int(W*0.75), int(H*0.03), int(W*0.22), int(H*0.08),
                     tag_text, font_size=11, color_hex=KT_WHITE, align=PP_ALIGN.RIGHT)

    line = slide.shapes.add_shape(1, 0, Emu(int(H*0.12)), Emu(W), Emu(int(H*0.005)))
    _fill(line, KT_RED)


def _footer(slide, project_name: str, W: int, H: int):
    _add_textbox(slide, int(W*0.05), int(H*0.94), int(W*0.9), int(H*0.05),
                 f"{project_name}  |  KT 수주전략  |  {date.today().strftime('%Y.%m.%d')}",
                 font_size=10, color_hex=KT_GRAY, align=PP_ALIGN.CENTER)


def _generate_slide_content(slide_data: dict, rfp_basics: dict, strategies: list[dict]) -> dict:
    """Claude API로 슬라이드 콘텐츠 생성"""
    client = anthropic.Anthropic()

    prompt = f"""슬라이드 콘텐츠를 생성하세요.

슬라이드 정보:
- 제목: {slide_data.get('title', '')}
- 핵심 메시지: {slide_data.get('key_message', '')}
- 내용 아웃라인: {json.dumps(slide_data.get('content_outline', []), ensure_ascii=False)}
- 평가 항목 태그: {slide_data.get('eval_tag', 'N/A')}
- 예상 점수 기여: {slide_data.get('score_target', 'N/A')}
- 논거 강도: {slide_data.get('argumentation', 'N/A')}

사업명: {rfp_basics.get('project_name', '')}
발주기관: {rfp_basics.get('client', '')}

순수 JSON만:
{{
  "headline": "슬라이드 핵심 한 줄 (임팩트 있게)",
  "bullets": ["핵심 불릿1", "핵심 불릿2", "핵심 불릿3"],
  "sub_bullets": {{"핵심 불릿1": ["세부 내용1", "세부 내용2"]}},
  "notes": "발표자 노트 (평가위원 설득 포인트)"
}}
"""

    response = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=1000,
        messages=[{"role": "user", "content": prompt}]
    )

    raw = response.content[0].text.strip()
    if raw.startswith("```"):
        raw = re.sub(r"^```(?:json)?\n?", "", raw)
        raw = re.sub(r"\n?```$", "", raw)
    return json.loads(raw)


def build_ppt(
    slides: list[dict],
    layer: int,
    rfp_basics: dict,
    strategies: list[dict],
    project_name: str,
    output_dir: Path,
) -> Path:
    """
    PPT 빌드.
    layer: 1, 2, 3 중 선택
    """
    layer_slides = get_layer_slides(slides, layer)
    prs = _build_prs()
    W = 9144000
    H = 5143500

    for slide_data in layer_slides:
        slide_type = slide_data.get("type", "content")

        if slide_type == "cover":
            sl = _blank(prs)
            bg = sl.shapes.add_shape(1, 0, 0, Emu(W), Emu(H))
            _fill(bg, KT_PRIMARY)
            accent = sl.shapes.add_shape(1, Emu(int(W*0.58)), 0, Emu(int(W*0.42)), Emu(H))
            _fill(accent, KT_ACCENT)
            bar = sl.shapes.add_shape(1, 0, Emu(int(H*0.55)), Emu(int(W*0.58)), Emu(int(H*0.008)))
            _fill(bar, KT_RED)
            _add_textbox(sl, int(W*0.06), int(H*0.28), int(W*0.50), int(H*0.22),
                         rfp_basics.get("project_name", slide_data.get("title", "")),
                         font_size=32, bold=True, color_hex=KT_WHITE, wrap=True)
            _add_textbox(sl, int(W*0.06), int(H*0.57), int(W*0.50), int(H*0.10),
                         f"For {rfp_basics.get('client', '')}",
                         font_size=18, color_hex="CCDDFF")
            _add_textbox(sl, int(W*0.06), int(H*0.70), int(W*0.50), int(H*0.08),
                         f"KT 주식회사  |  Layer {layer} 제안서",
                         font_size=13, color_hex="AABBCC")
        else:
            content = _generate_slide_content(slide_data, rfp_basics, strategies)
            sl = _blank(prs)
            bg = sl.shapes.add_shape(1, 0, 0, Emu(W), Emu(H))
            _fill(bg, KT_LIGHT)

            _header(
                sl,
                slide_data.get("title", content.get("headline", "")),
                W, H,
                score_tag=slide_data.get("eval_tag", ""),
                score_target=slide_data.get("score_target", ""),
            )

            headline = content.get("headline", "")
            if headline:
                hl = sl.shapes.add_shape(1, Emu(int(W*0.05)), Emu(int(H*0.13)),
                                          Emu(int(W*0.90)), Emu(int(H*0.07)))
                _fill(hl, "E8F4FD")
                _add_textbox(sl, int(W*0.06), int(H*0.135), int(W*0.88), int(H*0.065),
                             headline, font_size=14, bold=True, color_hex=KT_PRIMARY)

            y = int(H * 0.23)
            bullets = content.get("bullets", [])
            sub_bullets = content.get("sub_bullets", {})
            for bullet in bullets:
                if y + int(H*0.08) > int(H*0.92):
                    break
                icon = sl.shapes.add_shape(1, Emu(int(W*0.05)), Emu(y + int(H*0.015)),
                                            Emu(int(W*0.006)), Emu(int(H*0.04)))
                _fill(icon, KT_RED)
                _add_textbox(sl, int(W*0.065), y, int(W*0.87), int(H*0.075),
                             bullet, font_size=14, color_hex=KT_DARK)
                y += int(H*0.075)
                for sub in sub_bullets.get(bullet, []):
                    if y + int(H*0.06) > int(H*0.92):
                        break
                    _add_textbox(sl, int(W*0.08), y, int(W*0.85), int(H*0.055),
                                 f"• {sub}", font_size=12, color_hex=KT_GRAY)
                    y += int(H*0.055)

            notes_text = content.get("notes", "")
            if slide_data.get("eval_tag"):
                notes_text = f"[평가항목: {slide_data['eval_tag']} | 기대점수: {slide_data.get('score_target','?')} | 논거강도: {slide_data.get('argumentation','?')}]\n\n" + notes_text
            if notes_text:
                sl.notes_slide.notes_text_frame.text = notes_text

            _footer(sl, project_name, W, H)

    output_dir.mkdir(parents=True, exist_ok=True)
    layer_names = {1: "임원용_5장", 2: "의사결정용_12장", 3: "전체제출용"}
    filename = output_dir / f"{project_name}_Layer{layer}_{layer_names[layer]}.pptx"
    prs.save(str(filename))
    return filename
