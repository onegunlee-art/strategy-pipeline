"""
python-pptx 기반 PPT 빌더
config/ppt_style.yaml 스타일을 사용하여 슬라이드를 렌더링합니다.
"""
from __future__ import annotations

from pathlib import Path
from datetime import date

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

ROOT = Path(__file__).parent.parent


# ── 헬퍼 ──────────────────────────────────────────────────────────

def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_textbox(slide, left, top, width, height,
                 text: str, font_name: str, font_size: int,
                 bold=False, color_hex="1A1A2E",
                 align=PP_ALIGN.LEFT, wrap=True) -> None:
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)


def _fill_shape(shape, hex_str: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex_str)
    shape.line.fill.background()


# ── 슬라이드별 빌더 ───────────────────────────────────────────────

class SlideBuilder:
    def __init__(self, prs: Presentation, style: dict, fonts: dict, sizes: dict, colors: dict):
        self.prs = prs
        self.style = style
        self.fonts = fonts
        self.sizes = sizes
        self.colors = colors
        W = style["slide"]["width_emu"]
        H = style["slide"]["height_emu"]
        self.W = W
        self.H = H
        self.MARGIN = int(W * 0.05)
        self.CONTENT_TOP = int(H * 0.22)
        self.CONTENT_W = W - 2 * self.MARGIN
        self.CONTENT_H = H - self.CONTENT_TOP - int(H * 0.05)

    def _blank_slide(self):
        blank_layout = self.prs.slide_layouts[6]  # blank
        return self.prs.slides.add_slide(blank_layout)

    def _add_header(self, slide, title: str) -> None:
        W, H = self.W, self.H
        colors = self.colors
        # 헤더 바
        bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            0, 0, Emu(W), Emu(int(H * 0.12))
        )
        _fill_shape(bar, colors["primary"])

        # 슬라이드 제목
        _add_textbox(
            slide,
            left=int(W * 0.04), top=int(H * 0.025),
            width=int(W * 0.82), height=int(H * 0.1),
            text=title,
            font_name=self.fonts["title"],
            font_size=self.sizes["slide_title"],
            bold=True,
            color_hex=colors["white"],
        )

        # 하단 구분선
        line = slide.shapes.add_shape(
            1, 0, Emu(int(H * 0.12)), Emu(W), Emu(int(H * 0.006))
        )
        _fill_shape(line, colors["secondary"])

    def _add_footer(self, slide, project: dict) -> None:
        W, H = self.W, self.H
        client = project.get("client", "")
        dt = project.get("date", "")
        footer_text = f"{client}  |  AI 서비스 수주전략  |  {dt}"
        _add_textbox(
            slide,
            left=self.MARGIN, top=int(H * 0.93),
            width=self.CONTENT_W, height=int(H * 0.06),
            text=footer_text,
            font_name=self.fonts["body"],
            font_size=self.sizes["caption"],
            color_hex=self.colors["gray"],
            align=PP_ALIGN.CENTER,
        )

    def build_cover(self, project: dict) -> None:
        slide = self._blank_slide()
        W, H = self.W, self.H
        colors = self.colors

        # 전체 배경
        bg = slide.shapes.add_shape(1, 0, 0, Emu(W), Emu(H))
        _fill_shape(bg, colors["primary"])

        # 대각선 강조 블록
        accent = slide.shapes.add_shape(1, Emu(int(W * 0.55)), 0, Emu(int(W * 0.45)), Emu(H))
        _fill_shape(accent, colors["accent"])
        accent.fill.fore_color.rgb = _rgb(colors["accent"])

        # 레드 포인트 바
        bar = slide.shapes.add_shape(1, 0, Emu(int(H * 0.55)), Emu(int(W * 0.55)), Emu(int(H * 0.008)))
        _fill_shape(bar, colors["secondary"])

        # 제목
        _add_textbox(
            slide,
            left=int(W * 0.06), top=int(H * 0.28),
            width=int(W * 0.48), height=int(H * 0.22),
            text=project.get("opportunity", "AI 서비스 수주전략"),
            font_name=self.fonts["title"],
            font_size=42,
            bold=True,
            color_hex=colors["white"],
            wrap=True,
        )

        # 클라이언트
        _add_textbox(
            slide,
            left=int(W * 0.06), top=int(H * 0.57),
            width=int(W * 0.48), height=int(H * 0.12),
            text=f"For {project.get('client', 'KT')}",
            font_name=self.fonts["body"],
            font_size=self.sizes["section_title"],
            color_hex=colors["light_bg"],
        )

        # 날짜
        _add_textbox(
            slide,
            left=int(W * 0.06), top=int(H * 0.72),
            width=int(W * 0.48), height=int(H * 0.08),
            text=project.get("date", ""),
            font_name=self.fonts["body"],
            font_size=self.sizes["body"],
            color_hex=colors["light_bg"],
        )

    def build_content(self, slide_cfg: dict, content: dict, project: dict) -> None:
        slide = self._blank_slide()
        W, H = self.W, self.H
        colors = self.colors
        MARGIN = self.MARGIN

        bg = slide.shapes.add_shape(1, 0, 0, Emu(W), Emu(H))
        _fill_shape(bg, colors["light_bg"])

        self._add_header(slide, content.get("title", slide_cfg.get("title", "")))

        # 헤드라인
        headline = content.get("headline", "")
        if headline:
            hl_box = slide.shapes.add_shape(
                1,
                Emu(MARGIN), Emu(int(H * 0.13)),
                Emu(self.CONTENT_W), Emu(int(H * 0.07))
            )
            _fill_shape(hl_box, "E8F4FD")
            _add_textbox(
                slide,
                left=MARGIN + int(W * 0.01),
                top=int(H * 0.135),
                width=self.CONTENT_W - int(W * 0.02),
                height=int(H * 0.065),
                text=headline,
                font_name=self.fonts["body"],
                font_size=self.sizes["body"],
                bold=True,
                color_hex=colors["primary"],
            )

        # 불릿 콘텐츠
        bullets = content.get("bullets", [])
        sub_bullets = content.get("sub_bullets", {})
        content_top = int(H * 0.22)

        y = content_top
        bullet_h = int(H * 0.07)
        sub_h = int(H * 0.055)

        for bullet in bullets:
            if y + bullet_h > int(H * 0.92):
                break

            # 불릿 아이콘 바
            icon = slide.shapes.add_shape(
                1,
                Emu(MARGIN), Emu(y + int(bullet_h * 0.2)),
                Emu(int(W * 0.007)), Emu(int(bullet_h * 0.6))
            )
            _fill_shape(icon, colors["secondary"])

            _add_textbox(
                slide,
                left=MARGIN + int(W * 0.018),
                top=y,
                width=self.CONTENT_W - int(W * 0.02),
                height=bullet_h,
                text=bullet,
                font_name=self.fonts["body"],
                font_size=self.sizes["body"],
                bold=False,
                color_hex=colors["dark_text"],
            )
            y += bullet_h

            # 서브 불릿
            for sub in sub_bullets.get(bullet, []):
                if y + sub_h > int(H * 0.92):
                    break
                _add_textbox(
                    slide,
                    left=MARGIN + int(W * 0.04),
                    top=y,
                    width=self.CONTENT_W - int(W * 0.05),
                    height=sub_h,
                    text=f"• {sub}",
                    font_name=self.fonts["body"],
                    font_size=self.sizes["caption"],
                    color_hex=colors["gray"],
                )
                y += sub_h

        # 발표자 노트
        notes = content.get("notes", "")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        self._add_footer(slide, project)


# ── 메인 빌드 함수 ─────────────────────────────────────────────────

def build_ppt(cfg: dict, strategy_content: dict) -> Path:
    style = cfg.get("style", {})
    colors = style.get("colors", {})
    fonts = style.get("fonts", {})
    font_sizes = style.get("font_sizes", {})
    slide_width = style["slide"]["width_emu"]
    slide_height = style["slide"]["height_emu"]

    prs = Presentation()
    prs.slide_width = Emu(slide_width)
    prs.slide_height = Emu(slide_height)

    builder = SlideBuilder(prs, style, fonts, font_sizes, colors)
    project = cfg.get("project", {})
    slides_cfg = cfg.get("slides", [])

    for slide_cfg in slides_cfg:
        if not slide_cfg.get("enabled", True):
            continue
        sid = slide_cfg["id"]

        if sid == "cover":
            builder.build_cover(project)
        else:
            content = strategy_content.get(sid, {})
            if not content:
                content = {"title": slide_cfg["title"], "headline": "", "bullets": [], "notes": ""}
            content.setdefault("title", slide_cfg["title"])
            builder.build_content(slide_cfg, content, project)

    output_dir = ROOT / cfg.get("output_dir", "output")
    output_dir.mkdir(exist_ok=True)
    dt = project.get("date", date.today().strftime("%Y-%m-%d"))
    client = project.get("client", "client")
    filename = output_dir / f"{client}_AI_strategy_{dt}.pptx"
    prs.save(filename)
    return filename
